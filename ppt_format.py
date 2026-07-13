"""Update PowerPoint text while preserving template run formatting."""

from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.text.text import _Paragraph, _Run
from pptx.util import Pt

from ppt_missing import has_numeric_value, has_text_value, missing_label


def set_paragraph_text_preserve(paragraph: _Paragraph, text: str) -> None:
    """Replace paragraph text without resetting font, color, or size."""
    if paragraph.runs:
        paragraph.runs[0].text = text
        for run in paragraph.runs[1:]:
            run.text = ""
        return
    paragraph.text = text


def set_text_frame_preserve(text_frame, text: str) -> None:
    """Replace all text in a text frame, keeping the first run's formatting."""
    if not text_frame.paragraphs:
        text_frame.text = text
        return
    set_paragraph_text_preserve(text_frame.paragraphs[0], text)
    for paragraph in text_frame.paragraphs[1:]:
        for run in paragraph.runs:
            run.text = ""


def set_cell_text_preserve(cell, text: str) -> None:
    """Update a table cell value without changing template styling."""
    if cell.text_frame.paragraphs:
        set_text_frame_preserve(cell.text_frame, text)
    else:
        cell.text = text


def set_cell_value_or_missing(
    cell,
    value,
    *,
    missing: str,
    text: str | None = None,
) -> None:
    """Write a value or a [MISSING: ...] marker — never leave stale template text."""
    if text is not None and text != "":
        set_cell_text_preserve(cell, text)
    elif has_numeric_value(value) or has_text_value(value):
        set_cell_text_preserve(cell, str(value).strip())
    else:
        set_cell_text_preserve(cell, missing_label(missing))


def clear_text_frame_content(text_frame, *, keep_first_paragraph: bool = True) -> None:
    """Clear visible text but keep paragraph/run structure for styling."""
    if not text_frame.paragraphs:
        return
    start = 0 if keep_first_paragraph else 1
    for idx, paragraph in enumerate(text_frame.paragraphs):
        if idx < start:
            set_paragraph_text_preserve(paragraph, "")
        else:
            for run in paragraph.runs:
                run.text = ""


def clear_table_data_rows(table, *, header_rows: int = 1) -> None:
    """Blank table body cells below the header rows."""
    for row_idx in range(header_rows, len(table.rows)):
        for col_idx in range(len(table.columns)):
            set_cell_text_preserve(table.cell(row_idx, col_idx), "")


def _hex_to_rgb(hex_color: str) -> RGBColor:
    cleaned = hex_color.strip().lstrip("#").upper()
    if len(cleaned) == 8:
        cleaned = cleaned[2:]
    return RGBColor.from_string(cleaned)


def set_cell_fill_solid(cell, hex_color: str | None) -> None:
    """Apply a solid background fill to a table cell."""
    if not hex_color:
        return
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(hex_color)


def set_cell_font_style(
    cell,
    *,
    bold: bool | None = None,
    color_hex: str | None = None,
    size_pt: float | None = None,
) -> None:
    """Set font attributes on the first run of a table cell."""
    text_frame = cell.text_frame
    if not text_frame.paragraphs:
        text_frame.text = text_frame.text or ""
    paragraph = text_frame.paragraphs[0]
    if not paragraph.runs:
        run = paragraph.add_run()
        run.text = paragraph.text or cell.text or ""
    else:
        run = paragraph.runs[0]
    if bold is not None:
        run.font.bold = bold
    if color_hex:
        run.font.color.rgb = _hex_to_rgb(color_hex)
    if size_pt is not None:
        run.font.size = Pt(size_pt)


def style_table_cell(
    cell,
    text: str,
    *,
    fill_hex: str | None = None,
    font_bold: bool | None = None,
    font_color_hex: str | None = None,
    font_size_pt: float | None = None,
) -> None:
    """Write cell text and apply template-style fill/font in one step."""
    set_cell_text_preserve(cell, text)
    set_cell_fill_solid(cell, fill_hex)
    if font_bold is not None or font_color_hex or font_size_pt is not None:
        set_cell_font_style(
            cell,
            bold=font_bold,
            color_hex=font_color_hex,
            size_pt=font_size_pt,
        )


def merge_table_row(table, row_idx: int, col_start: int, col_end: int) -> None:
    """Merge cells on one row from col_start through col_end (inclusive)."""
    if col_end <= col_start:
        return
    if col_start < 0 or col_end >= len(table.columns):
        return
    table.cell(row_idx, col_start).merge(table.cell(row_idx, col_end))
