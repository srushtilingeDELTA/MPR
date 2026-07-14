"""People slide (PPT 7) from New GSE MPR Workings.xlsx → PEOPLE tab.

Only the first PEOPLE metrics table matters, plus these three Excel charts:
  1. Leadership Engagement
  2. Psychological Safety
  3. Accountability
All three are placed in a right-side stack. Leading Issues / Action Plans stay
empty editable text boxes.
"""

from __future__ import annotations

import io
import logging
import tempfile
from pathlib import Path

from openpyxl import load_workbook
from PIL import Image
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Emu

from gir_panels import clear_leading_action_narrative
from scorecard_screenshots import (
    _find_com_worksheet,
    _find_sheet_name,
    _open_excel_workbook,
    _png_from_pil,
    _range_address,
    _validate_capture,
    capture_range_png,
    place_picture_on_slide,
    resolve_sheet_name,
)

logger = logging.getLogger(__name__)

# Template table slot from GSE MPR template slide 7.
PEOPLE_TABLE_BOX = (362_427, 1_036_821, 7_426_936, 2_255_520)

# Right-side chart column — three equal stacked slots.
PEOPLE_CHART_LEFT = 7_976_668
PEOPLE_CHART_WIDTH = 3_657_600
PEOPLE_CHART_TOP = 1_036_821
PEOPLE_CHART_BOTTOM = 6_200_000
PEOPLE_CHART_GAP = 70_000

# Required charts, in top→bottom display order.
PEOPLE_CHART_SPECS = [
    {
        "key": "leadership",
        "title": "Leadership Engagement",
        "match": ["leadership engagement", "leadership"],
        "patterns": ["Leadership Engagement", "LEADERSHIP ENGAGEMENT"],
    },
    {
        "key": "psychological",
        "title": "Psychological Safety",
        "match": ["psychological safety", "psychological", "psych safety"],
        "patterns": ["Psychological Safety", "PSYCHOLOGICAL SAFETY"],
    },
    {
        "key": "accountability",
        "title": "Accountability",
        "match": ["accountability"],
        "patterns": ["Accountability", "ACCOUNTABILITY"],
    },
]

_TABLE_END_TOKENS = (
    "where i work, employees are held accountable",
    "employees are held accountable",
)


def people_chart_boxes(count: int = 3) -> list[tuple[int, int, int, int]]:
    """Evenly stack `count` chart boxes in the right column."""
    count = max(1, int(count))
    usable = PEOPLE_CHART_BOTTOM - PEOPLE_CHART_TOP - PEOPLE_CHART_GAP * (count - 1)
    height = max(900_000, usable // count)
    boxes = []
    top = PEOPLE_CHART_TOP
    for idx in range(count):
        h = height
        if idx == count - 1:
            h = max(height, PEOPLE_CHART_BOTTOM - top)
        boxes.append((PEOPLE_CHART_LEFT, top, PEOPLE_CHART_WIDTH, h))
        top += height + PEOPLE_CHART_GAP
    return boxes


def _cell_str(ws, row: int, col: int) -> str:
    val = ws.cell(row, col).value
    if val is None:
        return ""
    return str(val).strip()


def _find_people_header(ws, *, max_row: int = 40, max_col: int = 12) -> tuple[int, int] | None:
    """Find the first PEOPLE scorecard header (PEOPLE + MTD/YTD/Score)."""
    for row in range(1, max_row + 1):
        for col in range(1, max_col + 1):
            text = _cell_str(ws, row, col).casefold()
            if text != "people":
                continue
            row_blob = " ".join(_cell_str(ws, row, c).casefold() for c in range(col, min(col + 8, max_col + 1)))
            next_blob = " ".join(
                _cell_str(ws, row + 1, c).casefold() for c in range(col, min(col + 8, max_col + 1))
            )
            if "mtd" in row_blob or "ytd" in row_blob or "score" in row_blob:
                return row, col
            if "actual" in next_blob or "goal" in next_blob:
                return row, col
    for row in range(1, max_row + 1):
        for col in range(1, max_col + 1):
            if "leadership engagement" in _cell_str(ws, row, col).casefold():
                return max(1, row - 2), col
    return None


def _discover_first_people_table(workbook_bytes: bytes, sheet_name: str) -> tuple[int, int, int, int]:
    """Locate only the first PEOPLE MTD/YTD/Score table (ignore later tables/dumps)."""
    wb = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
    try:
        match = _find_sheet_name(list(wb.sheetnames), sheet_name) or sheet_name
        ws = wb[match]
        hdr = _find_people_header(ws)
        if hdr is None:
            raise ValueError(f"Could not find PEOPLE scorecard header on {sheet_name!r}")

        start_row, start_col = hdr
        end_col = start_col + 6
        end_row = start_row
        for row in range(start_row, min(int(ws.max_row or 40), start_row + 16) + 1):
            label = _cell_str(ws, row, start_col).casefold()
            row_has = any(_cell_str(ws, row, col) for col in range(start_col, end_col + 1))
            if not row_has:
                break
            end_row = row
            if any(token in label for token in _TABLE_END_TOKENS):
                break
            if row > start_row + 2 and label in {"people", "kpi", "yr_nb", "entity"}:
                end_row = row - 1
                break

        if end_row < start_row + 8:
            end_row = min(int(ws.max_row or 40), start_row + 9)

        logger.info(
            "PEOPLE first table %s!%s",
            match,
            _range_address(start_row, end_row, start_col, end_col),
        )
        return start_row, end_row, start_col, end_col
    finally:
        wb.close()


def _com_chart_title(chart_obj) -> str:
    for getter in (
        lambda: chart_obj.Chart.ChartTitle.Text,
        lambda: chart_obj.Chart.HasTitle and chart_obj.Chart.ChartTitle.Text,
        lambda: chart_obj.Name,
    ):
        try:
            text = str(getter() or "").strip()
            if text:
                return text
        except Exception:
            continue
    return ""


def _match_chart_spec(title: str) -> dict | None:
    lower = (title or "").casefold()
    for spec in PEOPLE_CHART_SPECS:
        for token in spec["match"]:
            if token in lower:
                return spec
    return None


def _export_named_people_charts(workbook_path: Path, sheet_name: str) -> dict[str, bytes]:
    """Export PEOPLE sheet charts keyed by leadership/psychological/accountability."""
    excel = None
    wb = None
    by_key: dict[str, bytes] = {}
    ordered_fallback: list[bytes] = []
    try:
        excel, wb = _open_excel_workbook(workbook_path)
        ws = _find_com_worksheet(wb, sheet_name)
        ws.Activate()
        count = int(ws.ChartObjects().Count)
        print(f">>> PEOPLE sheet has {count} Excel chart object(s)")
        scored: list[tuple[float, str, bytes]] = []
        for idx in range(1, count + 1):
            chart_obj = ws.ChartObjects(idx)
            export_path = Path(tempfile.gettempdir()) / f"mpr_people_chart_{idx}.png"
            try:
                if export_path.exists():
                    export_path.unlink()
                top = float(getattr(chart_obj, "Top", idx * 100))
                title = _com_chart_title(chart_obj)
                chart_obj.Chart.Export(str(export_path))
                if not (export_path.exists() and export_path.stat().st_size > 1500):
                    continue
                data = export_path.read_bytes()
                try:
                    data = _validate_capture(
                        data, label=f"PEOPLE chart {idx}", min_w=80, min_h=60, require_wide=False
                    )
                except Exception:
                    with Image.open(io.BytesIO(data)) as img:
                        data = _png_from_pil(img)
                scored.append((top, title, data))
                logger.info("Exported PEOPLE chart %s title=%r (%s bytes)", idx, title, len(data))
            except Exception as exc:
                logger.info("PEOPLE chart %s export failed: %s", idx, exc)
            finally:
                try:
                    if export_path.exists():
                        export_path.unlink()
                except Exception:
                    pass

        scored.sort(key=lambda item: item[0])
        for top, title, data in scored:
            spec = _match_chart_spec(title)
            if spec and spec["key"] not in by_key:
                by_key[spec["key"]] = data
                print(f">>> Matched PEOPLE chart '{spec['title']}' from Excel title {title!r}")
            else:
                ordered_fallback.append(data)

        # Fill any unmatched required keys from remaining top→bottom exports.
        for spec in PEOPLE_CHART_SPECS:
            if spec["key"] in by_key:
                continue
            if not ordered_fallback:
                break
            by_key[spec["key"]] = ordered_fallback.pop(0)
            print(f">>> Assigned unmatched PEOPLE chart export to '{spec['title']}' (position order)")
    finally:
        try:
            if wb is not None:
                wb.Close(SaveChanges=False)
        except Exception:
            pass
        try:
            if excel is not None:
                excel.Quit()
        except Exception:
            pass
    return by_key


def _remove_people_table_and_pictures(slide) -> int:
    """Remove PEOPLE table and prior content pictures; keep charts for native fallback."""
    removed = 0
    for shape in list(slide.shapes):
        drop = False
        if getattr(shape, "has_table", False):
            drop = True
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE and 900_000 <= int(shape.top) < 6_000_000:
            drop = True
        if not drop:
            continue
        shape._element.getparent().remove(shape._element)
        removed += 1
    return removed


def _remove_people_charts(slide) -> int:
    removed = 0
    for shape in list(slide.shapes):
        if getattr(shape, "has_chart", False):
            shape._element.getparent().remove(shape._element)
            removed += 1
    return removed


def _monthly_from_data(data, patterns: list[str]) -> tuple[list[float | None], float | None]:
    """Best-effort monthly series from Actuals (used when Excel chart export is incomplete)."""
    try:
        monthly = data.monthly_series(patterns, through_month=data.month, workbook="actuals")
        goal = data.kpi_value(patterns, month=data.month, workbook="actuals").goal
        return monthly, goal
    except Exception as exc:
        logger.info("Could not load Actuals series for %s: %s", patterns, exc)
        return [None] * 12, None


def _add_fallback_chart(slide, box: tuple[int, int, int, int], title: str, monthly, goal) -> bool:
    """Create a native PPT chart in the slot when Workings chart export is missing."""
    from ppt_builder import CHART_CATEGORIES, _safe_number

    if not any(_safe_number(v) is not None for v in (monthly or [])):
        return False
    left, top, width, height = box
    ytd_vals = [v for v in monthly[:12] if _safe_number(v) is not None]
    ytd = sum(ytd_vals) / len(ytd_vals) if ytd_vals else None
    chart_data = CategoryChartData()
    chart_data.categories = CHART_CATEGORIES
    chart_data.add_series("Actual", tuple(list(monthly[:12]) + [ytd]))
    if _safe_number(goal) is not None:
        chart_data.add_series("Goal", tuple([goal] * 13))
    shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(left),
        Emu(top),
        Emu(width),
        Emu(height),
        chart_data,
    )
    try:
        shape.chart.has_title = True
        shape.chart.chart_title.text_frame.text = title
    except Exception:
        try:
            shape.name = f"PeopleChart_{title}"
        except Exception:
            pass
    print(f">>> PEOPLE fallback native chart added for '{title}'")
    return True


def apply_people_workings_panels(slide, data, element: dict) -> bool:
    """Fill slide 7 from Workings!PEOPLE: first table + 3 named charts on the right."""
    workbook = element.get("workbook", "workings")
    prefer_com = bool(element.get("prefer_excel_com", True))
    fit = str(element.get("fit", "fill")).lower()
    chart_limit = int(element.get("chart_count", 3) or 3)
    allow_fallback = bool(element.get("chart_fallback", True))

    try:
        workbook_bytes = data.store.workbook_bytes(workbook)
    except FileNotFoundError as exc:
        logger.warning("PEOPLE workings screenshot skipped: %s", exc)
        return False

    available = []
    try:
        available = list(data.sheet_names(workbook))
    except Exception:
        available = []

    try:
        sheet_name = resolve_sheet_name(
            workbook_bytes,
            sheet=element.get("sheet", "PEOPLE"),
            sheet_index=element.get("sheet_index"),
            sheet_match=element.get("sheet_match", ["PEOPLE", "People"]),
            sheet_match_index=int(element.get("sheet_match_index", 0) or 0),
            available=available or None,
        )
    except Exception as exc:
        logger.warning("Could not resolve Workings PEOPLE sheet: %s", exc)
        return False

    out_dir = Path(getattr(data.store, "base_dir", Path("."))) / "output"
    out_dir.mkdir(parents=True, exist_ok=True)
    placed = 0

    # 1) First PEOPLE metrics table only.
    table_png = None
    try:
        start_row, end_row, start_col, end_col = _discover_first_people_table(workbook_bytes, sheet_name)
        table_png = capture_range_png(
            workbook_bytes,
            sheet_name,
            start_row,
            end_row,
            start_col,
            end_col,
            prefer_excel_com=prefer_com,
        )
        table_png = _validate_capture(
            table_png, label=f"PEOPLE table {sheet_name}", min_w=80, min_h=40, require_wide=False
        )
    except Exception as exc:
        logger.warning("PEOPLE first-table capture failed: %s", exc)
        table_png = None

    # 2) Export the three named charts from Workings!PEOPLE.
    charts_by_key: dict[str, bytes] = {}
    if prefer_com:
        try:
            with tempfile.TemporaryDirectory(prefix="mpr_people_") as tmp:
                path = Path(tmp) / "workings.xlsx"
                path.write_bytes(workbook_bytes)
                charts_by_key = _export_named_people_charts(path, sheet_name)
        except Exception as exc:
            logger.info("PEOPLE Excel chart export unavailable: %s", exc)

    if not table_png and not charts_by_key:
        logger.warning("No PEOPLE table/charts captured from workings/%s", sheet_name)
        return False

    removed = _remove_people_table_and_pictures(slide)
    logger.info("Slide 7 cleared %s table/picture shape(s)", removed)

    if table_png:
        left, top, width, height = PEOPLE_TABLE_BOX
        place_picture_on_slide(
            slide,
            table_png,
            left=left,
            top=top,
            max_width=width,
            max_height=height,
            fit=fit,
        )
        placed += 1
        try:
            with Image.open(io.BytesIO(table_png)) as img:
                debug = out_dir / f"_debug_people_table_{img.width}x{img.height}.png"
                img.save(debug)
                print(
                    f">>> PEOPLE first table placed from workings/{sheet_name} "
                    f"({img.width}x{img.height}) -> {debug.name}"
                )
        except Exception:
            print(f">>> PEOPLE first table placed from workings/{sheet_name}")

    boxes = people_chart_boxes(chart_limit)
    specs = PEOPLE_CHART_SPECS[:chart_limit]
    screenshot_count = sum(1 for spec in specs if spec["key"] in charts_by_key)

    # If we have Workings chart screenshots, replace native charts and place all three.
    if screenshot_count:
        _remove_people_charts(slide)
        for idx, (spec, box) in enumerate(zip(specs, boxes)):
            png = charts_by_key.get(spec["key"])
            if png is None:
                continue
            left, top, width, height = box
            place_picture_on_slide(
                slide,
                png,
                left=left,
                top=top,
                max_width=width,
                max_height=height,
                fit=fit,
            )
            placed += 1
            try:
                with Image.open(io.BytesIO(png)) as img:
                    debug = out_dir / f"_debug_people_{spec['key']}_{img.width}x{img.height}.png"
                    img.save(debug)
                    print(
                        f">>> PEOPLE chart {idx+1}/3 '{spec['title']}' placed "
                        f"({img.width}x{img.height}) -> {debug.name}"
                    )
            except Exception:
                print(f">>> PEOPLE chart {idx+1}/3 '{spec['title']}' placed")

    # Ensure all three right-side slots are filled (fallback native charts for gaps).
    missing = [spec for spec in specs if spec["key"] not in charts_by_key]
    if missing and allow_fallback:
        if screenshot_count == 0:
            # No exports at all — clear template charts and rebuild all three slots.
            _remove_people_charts(slide)
        for spec in missing:
            idx = next(i for i, s in enumerate(specs) if s["key"] == spec["key"])
            box = boxes[idx]
            monthly, goal = _monthly_from_data(data, spec["patterns"])
            if _add_fallback_chart(slide, box, spec["title"], monthly, goal):
                placed += 1
            else:
                logger.warning("Could not create fallback chart for %s", spec["title"])
                print(f">>> WARNING: missing PEOPLE chart '{spec['title']}'")
    elif missing:
        for spec in missing:
            print(f">>> WARNING: missing PEOPLE chart '{spec['title']}'")

    print(
        ">>> PEOPLE right-side charts required: "
        + ", ".join(spec["title"] for spec in specs)
    )

    if bool(element.get("clear_narrative", True)):
        n = clear_leading_action_narrative(slide)
        print(f">>> PEOPLE Leading Issues / Action Plans cleared ({n} text box(es))")

    print(
        f"\n>>> Slide 7 People: placed {placed} item(s) from workings/{sheet_name} "
        f"(table={'yes' if table_png else 'no'}, "
        f"chart screenshots={screenshot_count}/3, "
        f"fallback native={len(missing) if allow_fallback else 0})\n"
    )
    return placed > 0
