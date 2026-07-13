"""Fit system scorecard tables to template placeholder dimensions."""

from __future__ import annotations

import logging
from typing import TYPE_CHECKING

import pandas as pd
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu

from scorecard_data import MONTH_HEADERS, format_scorecard_cell
from scorecard_style import apply_scorecard_block_styles

if TYPE_CHECKING:
    pass

logger = logging.getLogger(__name__)

EMU_PER_INCH = 914400
MIN_FONT_PT = 5.0
MAX_FONT_PT = 9.0
DENSE_MARGIN_EMU = 36000
NORMAL_MARGIN_EMU = 45720


def compute_scorecard_font_size(nrows: int, height_emu: int) -> float:
    """Pick a font size so dense scorecard rows fit the placeholder height."""
    if nrows <= 0 or height_emu <= 0:
        return 8.0
    row_height_in = height_emu / EMU_PER_INCH / nrows
    size = row_height_in * 72 * 0.58
    return max(MIN_FONT_PT, min(MAX_FONT_PT, round(size, 1)))


def scorecard_column_weights(ncols: int) -> list[float]:
    """Relative column widths: narrow spacer/weight/month, wide KPI."""
    if ncols <= 0:
        return []
    if ncols == 1:
        return [1.0]
    weights = [0.25, 0.75, 2.6, 0.95]
    if ncols <= 4:
        return weights[:ncols]
    month_weight = 0.58 if ncols <= 18 else 0.5
    return weights + [month_weight] * (ncols - 4)


def _kpi_column_index(block: pd.DataFrame) -> int:
    for row_idx in range(min(8, len(block))):
        for col_idx in range(len(block.columns)):
            header = format_scorecard_cell(block.iat[row_idx, col_idx]).strip().upper()
            if header == "KPI":
                return col_idx
    return 2 if len(block.columns) > 2 else 0


def _numeric_column_indices(block: pd.DataFrame, kpi_col: int) -> set[int]:
    numeric_cols: set[int] = set()
    for row_idx in range(min(8, len(block))):
        for col_idx in range(len(block.columns)):
            header = format_scorecard_cell(block.iat[row_idx, col_idx]).strip().upper()
            if not header:
                continue
            if header in MONTH_HEADERS or header in {"WEIGHT", "TOTAL SCORE", "YTD", "YE"}:
                numeric_cols.add(col_idx)
            elif col_idx > kpi_col:
                numeric_cols.add(col_idx)
    if not numeric_cols:
        numeric_cols = {idx for idx in range(len(block.columns)) if idx != kpi_col and idx != 0}
    return numeric_cols


def apply_scorecard_table_layout(
    table,
    block: pd.DataFrame,
    *,
    width_emu: int,
    height_emu: int,
) -> float:
    """
    Distribute column widths and row heights to fill the placeholder area.
    Returns the font size chosen for the table body.
    """
    if block.empty or width_emu <= 0 or height_emu <= 0:
        return 8.0

    nrows = len(table.rows)
    ncols = len(table.columns)
    if nrows == 0 or ncols == 0:
        return 8.0

    font_size = compute_scorecard_font_size(nrows, height_emu)
    margin = DENSE_MARGIN_EMU if nrows > 24 else NORMAL_MARGIN_EMU
    row_height = max(1, height_emu // nrows)

    weights = scorecard_column_weights(ncols)
    total_weight = sum(weights)
    assigned_width = 0
    for col_idx, weight in enumerate(weights):
        if col_idx == ncols - 1:
            col_width = max(1, width_emu - assigned_width)
        else:
            col_width = max(1, int(width_emu * weight / total_weight))
            assigned_width += col_width
        table.columns[col_idx].width = Emu(col_width)

    assigned_height = 0
    for row_idx in range(nrows):
        if row_idx == nrows - 1:
            row_h = max(1, height_emu - assigned_height)
        else:
            row_h = row_height
            assigned_height += row_h
        table.rows[row_idx].height = Emu(row_h)

    kpi_col = min(_kpi_column_index(block), ncols - 1)
    numeric_cols = _numeric_column_indices(block, kpi_col)

    for row_idx in range(nrows):
        for col_idx in range(ncols):
            cell = table.cell(row_idx, col_idx)
            text_frame = cell.text_frame
            text_frame.margin_left = Emu(margin)
            text_frame.margin_right = Emu(margin)
            text_frame.margin_top = Emu(margin // 2)
            text_frame.margin_bottom = Emu(margin // 2)
            text_frame.word_wrap = col_idx == kpi_col
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

            if not text_frame.paragraphs:
                continue
            paragraph = text_frame.paragraphs[0]
            if col_idx == kpi_col:
                paragraph.alignment = PP_ALIGN.LEFT
            elif col_idx in numeric_cols or col_idx == 0:
                paragraph.alignment = PP_ALIGN.CENTER
            else:
                paragraph.alignment = PP_ALIGN.CENTER

    logger.debug(
        "Scorecard layout: %s rows x %s cols, font=%spt, area=%sx%s in",
        nrows,
        ncols,
        font_size,
        round(width_emu / EMU_PER_INCH, 2),
        round(height_emu / EMU_PER_INCH, 2),
    )
    return font_size


def fit_scorecard_table(
    table,
    block: pd.DataFrame,
    *,
    width_emu: int,
    height_emu: int,
) -> None:
    """Apply layout, colors, and scaled fonts so the table matches the template screenshot."""
    font_size = apply_scorecard_table_layout(table, block, width_emu=width_emu, height_emu=height_emu)
    apply_scorecard_block_styles(table, block, font_size_pt=font_size)
