"""
GSE MPR Report Generator — Phase 1
Manual Excel download -> filter report month -> PowerPoint from template
"""

from __future__ import annotations

import logging
from datetime import date
from pathlib import Path

import pandas as pd
import yaml
from dateutil.relativedelta import relativedelta
from pptx import Presentation

SCRIPT_VERSION = "2026.07.06-may"

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s",
)
logger = logging.getLogger(__name__)

BASE_DIR = Path(__file__).resolve().parent

MONTH_FROM_SHEET = {
    "jan": 1, "january": 1,
    "feb": 2, "february": 2,
    "mar": 3, "march": 3,
    "apr": 4, "april": 4,
    "may": 5,
    "jun": 6, "june": 6,
    "jul": 7, "july": 7,
    "aug": 8, "august": 8,
    "sep": 9, "sept": 9, "september": 9,
    "oct": 10, "october": 10,
    "nov": 11, "november": 11,
    "dec": 12, "december": 12,
}


def parse_bool(value, default: bool = True) -> bool:
    """Parse YAML/CLI-style booleans reliably."""
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, (int, float)):
        return value != 0
    if isinstance(value, str):
        normalized = value.strip().lower()
        if normalized in ("true", "yes", "1", "on"):
            return True
        if normalized in ("false", "no", "0", "off"):
            return False
    return default


def month_from_sheet_name(sheet_name: str) -> int | None:
    if not sheet_name:
        return None
    token = sheet_name.strip().split()[0].lower()
    return MONTH_FROM_SHEET.get(token)


def load_config(path: Path | None = None) -> dict:
    config_path = path or BASE_DIR / "config.yaml"
    if not config_path.exists():
        raise FileNotFoundError(
            f"Config not found: {config_path}\n"
            "Create config.yaml in the MPR folder (see README)."
        )

    with open(config_path, "r", encoding="utf-8") as f:
        config = yaml.safe_load(f) or {}

    excel_cfg = config.get("excel")
    if not isinstance(excel_cfg, dict):
        raise ValueError("config.yaml must include an 'excel:' section.")

    required_excel_keys = ("file_path", "sheet_name")
    missing_excel = [key for key in required_excel_keys if key not in excel_cfg]
    if missing_excel:
        raise ValueError(
            "config.yaml excel section is missing: "
            f"{missing_excel}\n"
            "Expected:\n"
            "  excel:\n"
            '    file_path: "data/MPR Actuals and Goals_v2.xlsx"\n'
            '    sheet_name: "May Actuals"\n'
            "    header_row: 0"
        )

    return config


def get_report_period(config: dict) -> tuple[int, int, str]:
    """Return (year, month, label) from config or previous calendar month."""
    report_cfg = config.get("report", {})
    excel_cfg = config.get("excel", {})
    use_previous = parse_bool(report_cfg.get("use_previous_month"), default=True)

    if use_previous:
        today = date.today()
        prev = today.replace(day=1) - relativedelta(months=1)
        return prev.year, prev.month, prev.strftime("%B %Y")

    year = report_cfg.get("year")
    month = report_cfg.get("month")

    if month is None:
        month = month_from_sheet_name(str(excel_cfg.get("sheet_name", "")))

    if year is None:
        year = date.today().year

    year = int(year)
    month = int(month)
    label = date(year, month, 1).strftime("%B %Y")
    return year, month, label


def load_excel(config: dict) -> pd.DataFrame:
    excel_cfg = config["excel"]
    file_path = BASE_DIR / excel_cfg["file_path"]

    if not file_path.exists():
        raise FileNotFoundError(
            f"Excel file not found: {file_path}\n"
            "Download from SharePoint and save it in the data/ folder."
        )

    logger.info("Reading Excel: %s", file_path)
    read_kwargs = {
        "sheet_name": excel_cfg["sheet_name"],
        "engine": "openpyxl",
    }
    if "header_row" in excel_cfg:
        read_kwargs["header"] = excel_cfg["header_row"]

    df = pd.read_excel(file_path, **read_kwargs)

    df = df.dropna(axis=1, how="all")
    df.columns = df.columns.astype(str).str.strip()
    logger.info("Loaded %s rows, %s columns", len(df), len(df.columns))
    return df


def filter_report_month(df: pd.DataFrame, config: dict) -> pd.DataFrame:
    cols = config["columns"]
    year, month, label = get_report_period(config)

    required = [cols["year"], cols["month"], cols["kpi"], cols["actual"]]
    missing = [c for c in required if c not in df.columns]
    if missing:
        raise ValueError(
            f"Missing columns: {missing}\nFound columns: {list(df.columns)}"
        )

    filtered = df.copy()
    filtered[cols["year"]] = pd.to_numeric(filtered[cols["year"]], errors="coerce")
    filtered[cols["month"]] = pd.to_numeric(filtered[cols["month"]], errors="coerce")
    filtered[cols["actual"]] = pd.to_numeric(filtered[cols["actual"]], errors="coerce")

    result = filtered.loc[
        (filtered[cols["year"]] == year) & (filtered[cols["month"]] == month)
    ].copy()

    logger.info("Report period: %s (%s rows)", label, len(result))
    return result


def summarize_by_region_kpi(df: pd.DataFrame, config: dict) -> pd.DataFrame:
    cols = config["columns"]
    group_cols = [cols["kpi"]]
    if cols["region"] in df.columns:
        group_cols.append(cols["region"])

    return (
        df.groupby(group_cols, dropna=False)[cols["actual"]]
        .sum()
        .reset_index()
        .sort_values(group_cols)
    )


def build_summary_table_text(summary: pd.DataFrame, config: dict) -> str:
    cols = config["columns"]
    lines = []

    for _, row in summary.head(25).iterrows():
        kpi = row[cols["kpi"]]
        region = row.get(cols["region"], "N/A")
        actual = row[cols["actual"]]
        actual_str = "N/A" if pd.isna(actual) else f"{actual:,.2f}"
        lines.append(f"{kpi} | {region} | {actual_str}")

    if len(summary) > 25:
        lines.append(f"... and {len(summary) - 25} more rows")

    return "\n".join(lines)


def fill_presentation(
    summary: pd.DataFrame,
    config: dict,
    report_label: str,
) -> Path:
    ppt_cfg = config["powerpoint"]
    template_path = BASE_DIR / ppt_cfg["template_path"]

    if not template_path.exists():
        raise FileNotFoundError(
            f"Template not found: {template_path}\n"
            "Add 'GSE MPR - Template.pptx' to the templates/ folder."
        )

    year, month, _ = get_report_period(config)
    output_name = ppt_cfg["output_name"].format(year=year, month=month)
    output_dir = BASE_DIR / ppt_cfg["output_dir"]
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / output_name

    prs = Presentation(str(template_path))
    table_text = build_summary_table_text(summary, config)

    if len(prs.slides) > 0:
        slide0 = prs.slides[0]
        if slide0.shapes.title:
            slide0.shapes.title.text = "GSE MPR Monthly Report"
        for shape in slide0.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = report_label
                break

    if len(prs.slides) > 1:
        slide1 = prs.slides[1]
        if slide1.shapes.title:
            slide1.shapes.title.text = f"KPI Summary — {report_label}"
        for shape in slide1.placeholders:
            if shape.placeholder_format.idx == 1:
                shape.text = table_text
                break
    else:
        layout = prs.slide_layouts[5] if len(prs.slide_layouts) > 5 else prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"KPI Summary — {report_label}"
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = table_text

    prs.save(output_path)
    logger.info("Saved report: %s", output_path)
    return output_path


def main() -> None:
    config = load_config()
    report_cfg = config.get("report", {})
    use_previous = parse_bool(report_cfg.get("use_previous_month"), default=True)
    year, month, report_label = get_report_period(config)
    logger.info(
        "Report settings: use_previous_month=%s, Yr_Nb=%s, Mo_Nb=%s (%s)",
        use_previous,
        year,
        month,
        report_label,
    )

    df = load_excel(config)
    monthly = filter_report_month(df, config)

    if monthly.empty:
        year, month, label = get_report_period(config)
        raise ValueError(
            f"No rows found for {label}. "
            f"Check Yr_Nb={year}, Mo_Nb={month} and re-download the workbook."
        )

    summary = summarize_by_region_kpi(monthly, config)
    output_path = fill_presentation(summary, config, report_label)
    print(f"\nDone. Report saved to:\n{output_path}")


if __name__ == "__main__":
    main()
    