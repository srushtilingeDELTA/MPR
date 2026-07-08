"""Print detailed template info for one slide (for slide-by-slide perfecting).

Usage:
    python scripts/slide_review.py 0
    python scripts/slide_review.py 4 --report   # also show what the builder targets
"""

from __future__ import annotations

import argparse
import sys
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR))

import io
import yaml
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from report_utils import load_config


def _load_prs(config: dict) -> Presentation:
    template_rel = config["powerpoint"]["template_path"]
    cached = config.get("_sharepoint_files", {}).get(template_rel)
    if cached:
        return Presentation(io.BytesIO(cached))
    path = BASE_DIR / template_rel
    return Presentation(str(path))


def _print_shape(shape, indent: int = 0) -> None:
    prefix = "  " * indent
    line = f"{prefix}- {shape.name!r} type={shape.shape_type}"
    if shape.has_text_frame and shape.text_frame.text.strip():
        text = shape.text_frame.text.replace("\n", " | ")[:160]
        line += f"\n{prefix}  text: {text!r}"
    if shape.has_table:
        table = shape.table
        headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
        line += f"\n{prefix}  table: {len(table.rows)}x{len(table.columns)} headers={headers}"
        if len(table.rows) > 1:
            sample = [table.cell(1, c).text.strip()[:20] for c in range(min(4, len(table.columns)))]
            line += f"\n{prefix}  row1 sample: {sample}"
    if shape.has_chart:
        title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else ""
        line += f"\n{prefix}  chart title: {title!r}"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        line += f"\n{prefix}  picture"
    print(line)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _print_shape(child, indent + 1)


def _map_entry_for_slide(slide_index: int) -> dict | None:
    path = BASE_DIR / "template_map.yaml"
    if not path.exists():
        return None
    with open(path, "r", encoding="utf-8") as handle:
        payload = yaml.safe_load(handle) or {}
    for slide in payload.get("slides", []):
        if slide.get("index") == slide_index:
            return slide
    return None


def main() -> int:
    parser = argparse.ArgumentParser(description="Review one template slide in detail.")
    parser.add_argument("slide", type=int, help="0-based slide index")
    parser.add_argument("--report", action="store_true", help="Show template_map.yaml entry")
    args = parser.parse_args()

    config = load_config(base_dir=BASE_DIR)
    template = BASE_DIR / config["powerpoint"]["template_path"]
    if not template.exists() and not config.get("_sharepoint_files"):
        print(f"Template not found: {template}\nSave template locally or run main.py first.")
        return 1

    prs = _load_prs(config)
    if args.slide < 0 or args.slide >= len(prs.slides):
        print(f"Slide {args.slide} out of range (0..{len(prs.slides) - 1})")
        return 1

    slide = prs.slides[args.slide]
    title = slide.shapes.title.text if slide.shapes.title else ""
    tables = sum(1 for s in slide.shapes if s.has_table)
    charts = sum(1 for s in slide.shapes if s.has_chart)
    pictures = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)

    print(f"SLIDE {args.slide} (user sees as slide {args.slide + 1} in PowerPoint)")
    print(f"  Title shape: {title!r}")
    print(f"  Tables: {tables}  Charts: {charts}  Pictures: {pictures}")
    print("  Shapes:")
    for shape in slide.shapes:
        _print_shape(shape, indent=2)

    if args.report:
        entry = _map_entry_for_slide(args.slide)
        print("\n  template_map.yaml:")
        if entry:
            print(yaml.dump(entry, default_flow_style=False, sort_keys=False))
        else:
            print("  (no entry — slide not mapped)")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())