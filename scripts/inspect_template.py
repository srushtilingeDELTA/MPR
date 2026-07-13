"""Inspect PowerPoint template: slides, tables, charts."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

BASE_DIR = Path(__file__).resolve().parent.parent


def walk_shapes(shapes, indent=0):
    for idx, shape in enumerate(shapes):
        prefix = "  " * indent
        line = f"{prefix}[{idx}] {shape.name} type={shape.shape_type}"
        if shape.has_text_frame and shape.text_frame.text.strip():
            text = shape.text_frame.text.replace("\n", " | ")[:120]
            line += f" text={text!r}"
        if shape.has_table:
            table = shape.table
            headers = [table.cell(0, c).text.strip() for c in range(len(table.columns))]
            line += f" TABLE {len(table.rows)}x{len(table.columns)} headers={headers}"
        if shape.has_chart:
            title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else ""
            line += f" CHART title={title!r}"
        print(line)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(shape.shapes, indent + 1)


def main() -> None:
    import yaml

    with open(BASE_DIR / "config.yaml", "r", encoding="utf-8") as f:
        config = yaml.safe_load(f)

    template = BASE_DIR / config["powerpoint"]["template_path"]
    prs = Presentation(str(template))
    print(f"Template: {template.name}")
    print(f"Slides: {len(prs.slides)}\n")

    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else ""
        tables = sum(1 for s in slide.shapes if s.has_table)
        charts = sum(1 for s in slide.shapes if s.has_chart)
        pics = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        if tables or charts or i in (0, 1):
            print(f"{'=' * 70}\nSLIDE {i} title={title!r} tables={tables} charts={charts} pictures={pics}")
            walk_shapes(slide.shapes)


if __name__ == "__main__":
    main()
