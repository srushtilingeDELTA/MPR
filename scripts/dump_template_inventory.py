"""Dump a detailed inventory of the PowerPoint template for template_map.yaml."""

from __future__ import annotations

import json
import sys
from pathlib import Path

BASE_DIR = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(BASE_DIR))

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from report_utils import load_config


def _shape_info(slide_index: int, shape, indent: int = 0) -> dict:
    info = {
        "index": indent,
        "name": shape.name,
        "type": str(shape.shape_type),
    }
    if shape.has_text_frame and shape.text_frame.text.strip():
        info["text"] = shape.text_frame.text.replace("\n", " | ")[:200]
    if shape.has_table:
        table = shape.table
        info["table"] = {
            "rows": len(table.rows),
            "cols": len(table.columns),
            "header": [table.cell(0, c).text.strip() for c in range(len(table.columns))],
        }
    if shape.has_chart:
        title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else ""
        info["chart"] = {"title": title}
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        info["children"] = [_shape_info(slide_index, child, indent + 1) for child in shape.shapes]
    return info


def main() -> int:
    config = load_config(base_dir=BASE_DIR)
    template_rel = config["powerpoint"]["template_path"]
    cached = config.get("_sharepoint_files", {}).get(template_rel)
    if cached:
        import io

        prs = Presentation(io.BytesIO(cached))
        source = "SharePoint cache"
    else:
        template = BASE_DIR / template_rel
        if not template.exists():
            print(f"Template not found: {template}")
            print("Run python main.py or sync SharePoint files first.")
            return 1
        prs = Presentation(str(template))
        source = str(template)

    inventory = {
        "source": source,
        "slide_count": len(prs.slides),
        "slides": [],
    }

    for i, slide in enumerate(prs.slides):
        title = slide.shapes.title.text if slide.shapes.title else ""
        slide_info = {
            "index": i,
            "title": title,
            "tables": sum(1 for s in slide.shapes if s.has_table),
            "charts": sum(1 for s in slide.shapes if s.has_chart),
            "pictures": sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE),
            "shapes": [_shape_info(i, shape) for shape in slide.shapes],
        }
        inventory["slides"].append(slide_info)

    out = BASE_DIR / "output" / "template_inventory.json"
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(inventory, indent=2), encoding="utf-8")
    print(f"Wrote {out}")
    print(f"Slides: {len(prs.slides)}")
    for slide in inventory["slides"]:
        if slide["tables"] or slide["charts"] or slide["pictures"]:
            print(
                f"  Slide {slide['index']:>2}: title={slide['title']!r} "
                f"tables={slide['tables']} charts={slide['charts']} pictures={slide['pictures']}"
            )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())