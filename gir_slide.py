"""Fill GIR slide tables and charts from MPR Actuals (+ scorecard scores)."""

from __future__ import annotations

from dataclasses import dataclass
from typing import TYPE_CHECKING

from data_lookup import lookup_kpi_exhaustive, lookup_monthly_series_exhaustive
from mpr_data import MONTH_LABELS, KpiValue
from pptx.enum.shapes import MSO_SHAPE_TYPE
from narrative_boxes import clear_manual_narrative_boxes
from ppt_format import set_cell_text_preserve
from ppt_missing import DNF, has_numeric_value
from scorecard_data import (
    build_scorecard_row_index,
    find_scorecard_header_row,
    format_scorecard_cell,
    load_system_scorecard,
)

if TYPE_CHECKING:
    from mpr_data import MprData

GIR_SCORECARD_KPI = "Global Injury Rate"
INJURY_BREAKDOWN_METRICS = {
    "total": ["Injury Count", "Injuries", "Injury", "Total Injuries"],
    "rec": ["Recordable", "Rec", "Recordable Injuries"],
    "nonrec": ["Non-Recordable", "NonRec", "Non Rec", "Non-Recordable Injuries"],
    "dart": ["DART", "Dart"],
}


def _iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)


@dataclass
class GirDataset:
    mtd: KpiValue
    monthly: list[float | None]
    ytd_actual: float | None
    injury_monthly: list[float | None]
    injury_ytd: float | None
    prior_year: list[float | None]
    score_mtd: float | None = None
    score_ytd: float | None = None


def _fmt_rate(value: float | None) -> str:
    if not has_numeric_value(value):
        return DNF
    return f"{float(value):.2f}"


def _fmt_count(value: float | None) -> str:
    if not has_numeric_value(value):
        return DNF
    return f"{round(float(value)):.0f}"


def _fmt_diff(actual: float | None, goal: float | None) -> str:
    if not has_numeric_value(actual) or not has_numeric_value(goal):
        return DNF
    diff = float(actual) - float(goal)
    if diff < 0:
        return f"({abs(diff):.2f})"
    return f"{diff:.2f}"


def _column_header_context(table, col_idx: int) -> str:
    parts: list[str] = []
    for row_idx in range(min(3, len(table.rows))):
        text = table.cell(row_idx, col_idx).text.strip().lower()
        if text:
            parts.append(text)
    for header_row in (0, 1):
        if header_row >= len(table.rows):
            continue
        for scan_col in range(col_idx, -1, -1):
            group = table.cell(header_row, scan_col).text.strip().lower()
            if group in {"mtd", "ytd", "score", "month to date", "year to date"}:
                parts.insert(0, group)
                break
        if parts and parts[0] in {"mtd", "ytd", "score", "month to date", "year to date"}:
            break
    return " | ".join(parts)


def parse_gir_metric_layout(table) -> dict[str, int]:
    layout: dict[str, int] = {}
    for col_idx in range(1, len(table.columns)):
        ctx = _column_header_context(table, col_idx)
        lower = ctx.lower()
        if "score" in lower and "ytd" in lower:
            layout["score_ytd"] = col_idx
        elif "score" in lower and ("mtd" in lower or "month" in lower):
            layout["score_mtd"] = col_idx
        elif ("ytd" in lower or "year to date" in lower) and "actual" in lower:
            layout["ytd_actual"] = col_idx
        elif ("ytd" in lower or "year to date" in lower) and "goal" in lower:
            layout["ytd_goal"] = col_idx
        elif ("ytd" in lower or "year to date" in lower) and ("+/-" in lower or "var" in lower or "vs" in lower):
            layout["ytd_var"] = col_idx
        elif ("mtd" in lower or "month to date" in lower) and "actual" in lower:
            layout["mtd_actual"] = col_idx
        elif ("mtd" in lower or "month to date" in lower) and "goal" in lower:
            layout["mtd_goal"] = col_idx
        elif ("mtd" in lower or "month to date" in lower) and ("+/-" in lower or "var" in lower or "vs" in lower):
            layout["mtd_var"] = col_idx
    return layout


def _scorecard_month_columns(block, header_row: int | None, data: MprData) -> tuple[int | None, int | None]:
    if header_row is None:
        return None, None
    month_label = MONTH_LABELS[data.month - 1].upper()
    month_label3 = month_label[:3]
    mtd_col = ytd_col = None
    for col_idx in range(len(block.columns)):
        header = format_scorecard_cell(block.iat[header_row, col_idx]).strip().upper()
        if header in {month_label, month_label3, data.report_month_short().upper()}:
            mtd_col = col_idx
        if header in {"YTD", "YE"}:
            ytd_col = col_idx
    return mtd_col, ytd_col


def _gir_scorecard_scores(data: MprData) -> tuple[float | None, float | None]:
    block = load_system_scorecard(data, block="summary_1")
    if block.empty:
        return None, None
    header_row = find_scorecard_header_row(block)
    mtd_col, ytd_col = _scorecard_month_columns(block, header_row, data)
    index = build_scorecard_row_index(block, header_row=header_row)
    row_idx = index.get(f"{GIR_SCORECARD_KPI.lower()}|score")
    if row_idx is None:
        return None, None
    mtd_score = ytd_score = None
    if mtd_col is not None:
        val = block.iat[row_idx, mtd_col]
        if format_scorecard_cell(val):
            mtd_score = float(val)
    if ytd_col is not None:
        val = block.iat[row_idx, ytd_col]
        if format_scorecard_cell(val):
            ytd_score = float(val)
    return mtd_score, ytd_score


def _lookup_mtd(data: MprData, patterns: list[str], workbook: str) -> KpiValue:
    value = data.kpi_value(
        patterns,
        month=data.month,
        workbook=workbook,
        system_only=True,
        aggregation="first",
    )
    if value.actual is not None or value.goal is not None:
        return value
    return lookup_kpi_exhaustive(data, patterns, month=data.month).value


def _lookup_monthly(data: MprData, patterns: list[str], workbook: str) -> list[float | None]:
    series = data.monthly_series(
        patterns,
        workbook=workbook,
        system_only=True,
        aggregation="first",
    )
    if any(v is not None for v in series):
        return series
    series, _ = lookup_monthly_series_exhaustive(data, patterns, aggregation="first")
    return series


def _lookup_ytd(data: MprData, patterns: list[str], workbook: str) -> float | None:
    ytd = data.ytd_value(
        patterns,
        workbook=workbook,
        system_only=True,
        aggregation="weighted",
    )
    if ytd is not None:
        return ytd
    return data.ytd_value(
        patterns,
        workbook=workbook,
        system_only=True,
        aggregation="first",
    )


def load_gir_dataset(data: MprData, config: dict, workbook: str) -> GirDataset:
    kpi = config.get("kpi_mappings", {})
    gir_patterns = [kpi.get("gir", "GIR")]
    injury_patterns = [kpi.get("injury_count", "Injury Count"), "Injuries", "Injury"]
    score_mtd, score_ytd = _gir_scorecard_scores(data)
    mtd = _lookup_mtd(data, gir_patterns, workbook)
    return GirDataset(
        mtd=mtd,
        monthly=_lookup_monthly(data, gir_patterns, workbook),
        ytd_actual=_lookup_ytd(data, gir_patterns, workbook),
        injury_monthly=data.monthly_series(
            injury_patterns,
            through_month=data.month,
            workbook=workbook,
            aggregation="sum",
            system_only=False,
        ),
        injury_ytd=data.ytd_sum(injury_patterns, workbook=workbook, system_only=False),
        prior_year=data.prior_year_monthly_series(
            gir_patterns,
            through_month=data.month,
            workbook=workbook,
        ),
        score_mtd=score_mtd,
        score_ytd=score_ytd,
    )


def _table_row_label(table, row_idx: int) -> str:
    return " ".join(table.cell(row_idx, col).text.strip().lower() for col in range(len(table.columns)))


def _find_table_row(table, *patterns: str) -> int | None:
    for row_idx in range(len(table.rows)):
        label = _table_row_label(table, row_idx)
        if all(pattern.lower() in label for pattern in patterns):
            return row_idx
    return None


def _clear_numeric_body(table, *, header_rows: int = 1, label_cols: int = 1) -> None:
    for row_idx in range(header_rows, len(table.rows)):
        for col_idx in range(label_cols, len(table.columns)):
            set_cell_text_preserve(table.cell(row_idx, col_idx), "")


def _fill_metric_table(table, dataset: GirDataset) -> None:
    gir_row = _find_table_row(table, "gir")
    if gir_row is None and len(table.rows) > 2:
        gir_row = 2
    if gir_row is None or table.cell(0, 0).text.strip().lower() != "metric":
        return

    layout = parse_gir_metric_layout(table)
    goal = dataset.mtd.goal
    _clear_numeric_body(table, header_rows=1, label_cols=2)

    defaults = {
        "mtd_actual": 1,
        "mtd_goal": 2,
        "mtd_var": 3,
        "ytd_actual": 4,
        "ytd_goal": 5,
        "ytd_var": 6,
        "score_mtd": 7,
        "score_ytd": 8,
    }
    cols = {key: layout.get(key, default) for key, default in defaults.items() if default < len(table.columns)}

    if "mtd_actual" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["mtd_actual"]), _fmt_rate(dataset.mtd.actual))
    if "mtd_goal" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["mtd_goal"]), _fmt_rate(goal))
    if "mtd_var" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["mtd_var"]), _fmt_diff(dataset.mtd.actual, goal))
    if "ytd_actual" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["ytd_actual"]), _fmt_rate(dataset.ytd_actual))
    if "ytd_goal" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["ytd_goal"]), _fmt_rate(goal))
    if "ytd_var" in cols:
        set_cell_text_preserve(table.cell(gir_row, cols["ytd_var"]), _fmt_diff(dataset.ytd_actual, goal))
    if "score_mtd" in cols:
        set_cell_text_preserve(
            table.cell(gir_row, cols["score_mtd"]),
            _fmt_rate(dataset.score_mtd) if dataset.score_mtd is not None else DNF,
        )
    if "score_ytd" in cols:
        set_cell_text_preserve(
            table.cell(gir_row, cols["score_ytd"]),
            _fmt_rate(dataset.score_ytd) if dataset.score_ytd is not None else DNF,
        )


def _fill_mtd_summary_table(table, data: MprData, dataset: GirDataset) -> None:
    header = table.cell(0, 0).text.strip().lower()
    if "actual:" not in _table_row_label(table, 1) and data.report_month_short().lower() not in header:
        return
    _clear_numeric_body(table, header_rows=1, label_cols=1)
    set_cell_text_preserve(table.cell(0, 0), data.report_month_short())
    set_cell_text_preserve(table.cell(1, 1), _fmt_rate(dataset.mtd.actual))
    if len(table.rows) > 2:
        set_cell_text_preserve(table.cell(2, 1), _fmt_rate(dataset.mtd.goal))


def _fill_yoy_table(table, data: MprData, gir_patterns: list[str], workbook: str) -> None:
    text = " ".join(table.cell(row, col).text.lower() for row in range(len(table.rows)) for col in range(len(table.columns)))
    if "yo1y" not in text and "yoy" not in text:
        return

    yo1y = data.kpi_value(gir_patterns, month=data.month, workbook=workbook, year=data.year - 1, system_only=True)
    yo2y = data.kpi_value(gir_patterns, month=data.month, workbook=workbook, year=data.year - 2, system_only=True)
    _clear_numeric_body(table, header_rows=1, label_cols=1)

    for row_idx in range(len(table.rows)):
        for col_idx in range(len(table.columns)):
            cell_text = table.cell(row_idx, col_idx).text.strip().lower()
            if "yo1y" in cell_text:
                target_col = col_idx + 1 if col_idx + 1 < len(table.columns) else col_idx
                set_cell_text_preserve(table.cell(row_idx, target_col), _fmt_rate(yo1y.actual))
            elif "yo2y" in cell_text:
                target_col = col_idx + 1 if col_idx + 1 < len(table.columns) else col_idx
                set_cell_text_preserve(table.cell(row_idx, target_col), _fmt_rate(yo2y.actual))


def _fill_recordable_table(table, data: MprData, dataset: GirDataset) -> None:
    if table.cell(0, 0).text.strip().lower() != "recordable":
        return

    gir_row = _find_table_row(table, "system", "gir") or _find_table_row(table, "gir") or 1
    injury_row = _find_table_row(table, "injury") or 2
    goal = dataset.mtd.goal
    _clear_numeric_body(table, header_rows=1, label_cols=1)

    for col_idx, month_num in enumerate(range(1, 13), start=1):
        if col_idx >= len(table.columns):
            break
        if month_num <= data.month:
            set_cell_text_preserve(table.cell(gir_row, col_idx), _fmt_rate(dataset.monthly[month_num - 1]))
            if injury_row is not None:
                set_cell_text_preserve(
                    table.cell(injury_row, col_idx),
                    _fmt_count(dataset.injury_monthly[month_num - 1]),
                )
        else:
            set_cell_text_preserve(table.cell(gir_row, col_idx), "")
            if injury_row is not None:
                set_cell_text_preserve(table.cell(injury_row, col_idx), "")

    if len(table.columns) > 13:
        set_cell_text_preserve(table.cell(gir_row, 13), _fmt_rate(dataset.ytd_actual))
        set_cell_text_preserve(table.cell(gir_row, 14), _fmt_diff(dataset.ytd_actual, goal))
    if injury_row is not None and len(table.columns) > 13:
        set_cell_text_preserve(table.cell(injury_row, 13), _fmt_count(dataset.injury_ytd))


def _fill_injury_breakdown_table(table, data: MprData, workbook: str) -> None:
    text = " ".join(table.cell(row, col).text.lower() for row in range(len(table.rows)) for col in range(len(table.columns)))
    if "dart" not in text or "nonrec" not in text:
        return

    col_map: dict[str, int] = {}
    for col_idx in range(len(table.columns)):
        header = table.cell(0, col_idx).text.strip().lower()
        if header in INJURY_BREAKDOWN_METRICS:
            col_map[header] = col_idx

    for row_idx in range(1, len(table.rows)):
        label = _table_row_label(table, row_idx)
        year: int | None = None
        if "2025" in label:
            year = data.year - 1
        elif "2026" in label:
            year = data.year
        elif "total" in label and "202" not in label:
            year = None
        else:
            continue

        entity = "587" if "587" in label else "613" if "613" in label else None

        for metric, patterns in INJURY_BREAKDOWN_METRICS.items():
            if metric not in col_map:
                continue
            if year is None:
                val = 0.0
                found = False
                for yr in (data.year - 1, data.year):
                    part = data.ytd_sum(
                        patterns,
                        workbook=workbook,
                        year=yr,
                        entity_pattern=entity,
                        system_only=False,
                    )
                    if part is not None:
                        val += part
                        found = True
                set_cell_text_preserve(
                    table.cell(row_idx, col_map[metric]),
                    _fmt_count(val) if found else DNF,
                )
            else:
                val = data.ytd_sum(
                    patterns,
                    workbook=workbook,
                    year=year,
                    entity_pattern=entity,
                    system_only=False,
                )
                set_cell_text_preserve(table.cell(row_idx, col_map[metric]), _fmt_count(val))


def fill_gir_tables(slide, data: MprData, config: dict, workbook: str) -> GirDataset:
    """Fill all GIR tables on the slide; return the loaded dataset for charts."""
    kpi = config.get("kpi_mappings", {})
    gir_patterns = [kpi.get("gir", "GIR")]
    dataset = load_gir_dataset(data, config, workbook)

    for shape in _iter_all_shapes(slide.shapes):
        if not shape.has_table:
            continue
        table = shape.table
        _fill_metric_table(table, dataset)
        _fill_mtd_summary_table(table, data, dataset)
        _fill_yoy_table(table, data, gir_patterns, workbook)
        _fill_recordable_table(table, data, dataset)
        _fill_injury_breakdown_table(table, data, workbook)

    return dataset


def clear_gir_narrative_boxes(slide) -> None:
    """Leave Leading Issues / Action Plan as blank textboxes (headers only)."""
    clear_manual_narrative_boxes(slide)
