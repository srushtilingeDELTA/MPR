"""Replace template screenshot pictures / OLE objects with native tables from Excel."""

from __future__ import annotations

import logging
import math
from typing import TYPE_CHECKING

import pandas as pd
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from ppt_format import set_cell_text_preserve
from scorecard_data import format_scorecard_cell
from scorecard_layout import fit_scorecard_table
from scorecard_style import apply_scorecard_block_styles, apply_worksheet_table_styles

if TYPE_CHECKING:
    from pptx.slide import Slide

logger = logging.getLogger(__name__)

MIN_DATA_PICTURE_WIDTH = Inches(4)
MIN_DATA_PICTURE_HEIGHT = Inches(2)
MIN_OLE_WIDTH = Inches(4)


def _shape_area(shape) -> int:
    return int(shape.width * shape.height)


def is_data_screenshot_shape(shape) -> bool:
    """True for large body pictures or embedded Excel/screenshot OLE objects."""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return shape.width >= MIN_DATA_PICTURE_WIDTH and shape.height >= MIN_DATA_PICTURE_HEIGHT
    if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT:
        return shape.width >= MIN_OLE_WIDTH
    return False


def find_largest_data_placeholder(slide: Slide):
    """Return the largest picture/OLE on the slide that holds screenshot data."""
    candidates = [s for s in slide.shapes if is_data_screenshot_shape(s)]
    if not candidates:
        return None
    return max(candidates, key=_shape_area)


def remove_data_placeholders(slide: Slide) -> int:
    """Remove screenshot pictures and embedded OLE data objects from the slide body."""
    removed = 0
    for shape in list(slide.shapes):
        if is_data_screenshot_shape(shape):
            shape._element.getparent().remove(shape._element)
            removed += 1
    if removed:
        logger.info("Removed %s data placeholder(s) (picture/OLE) from slide", removed)
    return removed


def _cell_text(value) -> str:
    if value is None or (isinstance(value, float) and (math.isnan(value) or math.isinf(value))):
        return ""
    if isinstance(value, (int, float)):
        return format_scorecard_cell(value)
    text = str(value).strip()
    return "" if text.lower() in {"nan", "none"} else text


def fill_table_from_dataframe(table, df: pd.DataFrame, *, start_row: int = 0, start_col: int = 0) -> int:
    """Write DataFrame values into an existing table; return cells written."""
    if df.empty:
        return 0
    filled = 0
    rows = min(len(table.rows) - start_row, len(df.index))
    cols = min(len(table.columns) - start_col, len(df.columns))
    for r in range(rows):
        for c in range(cols):
            text = _cell_text(df.iat[r, c])
            set_cell_text_preserve(table.cell(start_row + r, start_col + c), text)
            if text:
                filled += 1
    return filled


def fill_styled_scorecard_table(
    table,
    df: pd.DataFrame,
    *,
    start_row: int = 0,
    start_col: int = 0,
    width_emu: int | None = None,
    height_emu: int | None = None,
) -> int:
    """Write scorecard values and apply template section/header colors."""
    filled = fill_table_from_dataframe(table, df, start_row=start_row, start_col=start_col)
    if width_emu and height_emu:
        fit_scorecard_table(table, df, width_emu=width_emu, height_emu=height_emu)
    else:
        apply_scorecard_block_styles(table, df)
    return filled


def fill_styled_worksheet_table(
    table,
    df: pd.DataFrame,
    ws,
    *,
    origin_row: int = 0,
    origin_col: int = 0,
    start_row: int = 0,
    start_col: int = 0,
) -> int:
    """Write worksheet values and mirror Excel/template styling."""
    filled = fill_table_from_dataframe(table, df, start_row=start_row, start_col=start_col)
    apply_worksheet_table_styles(table, ws, df, origin_row=origin_row, origin_col=origin_col)
    return filled


def replace_placeholder_with_table(slide: Slide, placeholder, nrows: int, ncols: int):
    """Remove a picture/OLE and add a native table at the same position."""
    left, top, width, height = placeholder.left, placeholder.top, placeholder.width, placeholder.height
    placeholder._element.getparent().remove(placeholder._element)
    nrows = max(1, nrows)
    ncols = max(1, ncols)
    shape = slide.shapes.add_table(nrows, ncols, left, top, width, height)
    return shape.table


def paste_dataframe_at_placeholder(
    slide: Slide,
    df: pd.DataFrame,
    *,
    placeholder=None,
    style: str = "plain",
    ws=None,
    origin_row: int = 0,
    origin_col: int = 0,
) -> tuple[object | None, int]:
    """
    Replace the largest data placeholder with a table filled from df.
    Returns (table, cells_written) or (None, 0).

    style: plain | scorecard | worksheet
    """
    if df.empty:
        return None, 0
    target = placeholder or find_largest_data_placeholder(slide)
    if target is None:
        return None, 0
    width_emu, height_emu = target.width, target.height
    table = replace_placeholder_with_table(slide, target, len(df.index), len(df.columns))
    if style == "scorecard":
        filled = fill_styled_scorecard_table(table, df, width_emu=width_emu, height_emu=height_emu)
    elif style == "worksheet" and ws is not None:
        filled = fill_styled_worksheet_table(table, df, ws, origin_row=origin_row, origin_col=origin_col)
    else:
        filled = fill_table_from_dataframe(table, df)
    logger.info("Pasted %s-row dataframe into native table (%s cells, style=%s)", len(df.index), filled, style)
    return table, filled


def paste_dataframe_on_slide(
    slide: Slide,
    df: pd.DataFrame,
    *,
    style: str = "plain",
    ws=None,
    origin_row: int = 0,
    origin_col: int = 0,
) -> tuple[object | None, int]:
    """Remove data screenshots/OLE, then paste df at the largest removed region or existing placeholder."""
    placeholder = find_largest_data_placeholder(slide)
    if placeholder is None:
        return None, 0
    return paste_dataframe_at_placeholder(
        slide,
        df,
        placeholder=placeholder,
        style=style,
        ws=ws,
        origin_row=origin_row,
        origin_col=origin_col,
    )
