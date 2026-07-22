"""Capture System scorecard sections from Excel and prepare slide images.

Primary path (Windows + Excel): CopyPicture via COM for true screenshots.
Fallback (any OS): render openpyxl cell fills/text to PNG with Pillow.
"""

from __future__ import annotations

import io
import logging
import os
import tempfile
import time
from dataclasses import dataclass
from pathlib import Path

from openpyxl import load_workbook
from openpyxl.cell.cell import Cell
from openpyxl.styles.colors import COLOR_INDEX
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.worksheet import Worksheet
from PIL import Image, ImageDraw, ImageFont

logger = logging.getLogger(__name__)

# Content area under the System Scorecard title on slides 3/4 (EMU).
# Matches the template picture slot used in the live GSE deck.
DEFAULT_LEFT = 314628
DEFAULT_TOP = 997580
DEFAULT_WIDTH = 11534806
DEFAULT_HEIGHT = 4862840


def _norm_sheet_key(name: str) -> str:
    """Casefold + strip so trailing spaces in Excel sheet names still match."""
    return str(name or "").strip().casefold()


def _find_sheet_name(names: list[str], sheet_name: str) -> str | None:
    """Resolve an exact sheet name, ignoring leading/trailing whitespace."""
    key = _norm_sheet_key(sheet_name)
    for name in names:
        if _norm_sheet_key(name) == key:
            return name
    return None


@dataclass
class Section:
    """One category block on the System scorecard (e.g. Safety & Security)."""

    index: int
    title: str
    start_row: int  # 1-based inclusive (category body start)
    end_row: int  # 1-based inclusive
    start_col: int  # 1-based inclusive
    end_col: int  # 1-based inclusive
    is_black: bool = False

    @property
    def range_address(self) -> str:
        return (
            f"{get_column_letter(self.start_col)}{self.start_row}:"
            f"{get_column_letter(self.end_col)}{self.end_row}"
        )


@dataclass
class SystemLayout:
    """Detected System sheet layout: shared month header + category sections."""

    header_start_row: int
    header_end_row: int
    start_col: int
    end_col: int
    sections: list[Section]

    def capture_bounds_for(self, sections: list[Section]) -> tuple[int, int, int, int]:
        """Body bounds for selected sections (header handled separately when needed)."""
        if not sections:
            return self.header_start_row, self.header_end_row, self.start_col, self.end_col
        return (
            min(s.start_row for s in sections),
            max(s.end_row for s in sections),
            self.start_col,
            self.end_col,
        )

    def needs_header_stitch(self, sections: list[Section]) -> bool:
        """True when selected body does not sit directly under the month header."""
        if not sections:
            return False
        body_start = min(s.start_row for s in sections)
        return body_start > self.header_end_row + 1


MONTH_HEADER_TOKENS = {
    "JAN",
    "FEB",
    "MAR",
    "APR",
    "MAY",
    "JUN",
    "JUL",
    "AUG",
    "SEP",
    "SEPT",
    "OCT",
    "NOV",
    "DEC",
    "YTD",
    "YE",
    "JANUARY",
    "FEBRUARY",
    "MARCH",
    "APRIL",
    "JUNE",
    "JULY",
    "AUGUST",
    "SEPTEMBER",
    "OCTOBER",
    "NOVEMBER",
    "DECEMBER",
}

# Ordered System tab categories used when left-rail labels are shapes/images
# (openpyxl cannot always read those cells, but TOTAL SCORE rows still mark boundaries).
DEFAULT_SYSTEM_CATEGORIES = [
    "Safety & Security",
    "Customer Experience",
    "Operations",
    "People",
    "Finance",
    "Overall and Opportunities",
]

# KPI names that must NEVER be treated as category section titles.
KNOWN_KPI_NAMES = {
    "budget",
    "total hours",
    "hours",
    "ot",
    "psych. safety",
    "psych safety",
    "psychological safety",
    "global injury rate",
    "gir",
    "ea compliance",
    "asap",
    "in-service rate",
    "vos (stationary)",
    "vos",
    "pmi",
    "critical jam",
    "clear time",
    "mbr",
    "engagement",
    "retention",
    "training",
}

METRIC_LABELS = {
    "plan",
    "actual",
    "percent",
    "points",
    "score",
    "weight",
    "kpi",
    "total score",
}

CATEGORY_KEYWORDS = (
    "safety & security",
    "customer experience",
    "operations",
    "people",
    "finance",
    "financial",
    "overall",
    "opportunities",
    "opportunity",
)


def _theme_rgb(theme: int | None, tint: float | None = None) -> tuple[int, int, int] | None:
    # Approximate Office theme colors used by Delta scorecards.
    theme_map = {
        0: (255, 255, 255),
        1: (0, 0, 0),
        2: (238, 238, 238),
        3: (31, 73, 125),
        4: (79, 129, 189),
        5: (192, 80, 77),
        6: (155, 187, 89),
        7: (128, 100, 162),
        8: (75, 172, 198),
        9: (247, 150, 70),
    }
    if theme is None:
        return None
    base = theme_map.get(theme)
    if base is None:
        return (80, 80, 80)
    if not tint:
        return base
    r, g, b = base
    if tint < 0:
        factor = 1.0 + tint
        return (int(r * factor), int(g * factor), int(b * factor))
    return (
        int(r + (255 - r) * tint),
        int(g + (255 - g) * tint),
        int(b + (255 - b) * tint),
    )


def _color_rgb(color) -> tuple[int, int, int] | None:
    if color is None or color.type is None:
        return None
    if color.type == "rgb" and color.rgb:
        value = str(color.rgb)
        if value.startswith("urn:"):
            return None
        if len(value) == 8:
            value = value[2:]
        if len(value) == 6:
            try:
                return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
            except ValueError:
                return None
    if color.type == "theme":
        return _theme_rgb(color.theme, getattr(color, "tint", None))
    if color.type == "indexed" and color.indexed is not None:
        idx = int(color.indexed)
        if 0 <= idx < len(COLOR_INDEX):
            value = COLOR_INDEX[idx]
            if isinstance(value, str) and len(value) >= 6:
                value = value[-6:]
                try:
                    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))  # type: ignore[return-value]
                except ValueError:
                    return None
    return None


def _fill_rgb(cell: Cell) -> tuple[int, int, int] | None:
    fill = cell.fill
    if fill is None or fill.fill_type in (None, "none"):
        return None
    return _color_rgb(fill.fgColor) or _color_rgb(fill.bgColor)


def _luminance(rgb: tuple[int, int, int]) -> float:
    r, g, b = rgb
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def _is_dark(rgb: tuple[int, int, int] | None) -> bool:
    return rgb is not None and _luminance(rgb) < 55


def _cell_text(cell: Cell) -> str:
    value = cell.value
    if value is None:
        return ""
    return str(value).strip()


def _normalize_label(text: str) -> str:
    """Collapse NBSP / newlines / multi-space so TOTAL SCORE matches reliably."""
    if not text:
        return ""
    cleaned = (
        str(text)
        .replace("\xa0", " ")
        .replace("\u2007", " ")
        .replace("\u202f", " ")
        .replace("\r", " ")
        .replace("\n", " ")
        .replace("\t", " ")
    )
    return " ".join(cleaned.split()).casefold()


def _is_total_score_label(text: str) -> bool:
    label = _normalize_label(text)
    if not label:
        return False
    compact = label.replace(" ", "")
    return (
        label == "total score"
        or compact == "totalscore"
        or label.startswith("total score")
        or "total score" in label
    )


def _row_has_content(ws: Worksheet, row: int, min_col: int, max_col: int) -> bool:
    for col in range(min_col, max_col + 1):
        if _cell_text(ws.cell(row, col)):
            return True
        if _fill_rgb(ws.cell(row, col)) is not None:
            return True
    return False


def _used_bounds(ws: Worksheet) -> tuple[int, int, int, int]:
    min_row, max_row = None, None
    min_col, max_col = None, None
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is None and _fill_rgb(cell) is None:
                continue
            r, c = cell.row, cell.column
            min_row = r if min_row is None else min(min_row, r)
            max_row = r if max_row is None else max(max_row, r)
            min_col = c if min_col is None else min(min_col, c)
            max_col = c if max_col is None else max(max_col, c)
    for merged in ws.merged_cells.ranges:
        min_row = merged.min_row if min_row is None else min(min_row, merged.min_row)
        max_row = merged.max_row if max_row is None else max(max_row, merged.max_row)
        min_col = merged.min_col if min_col is None else min(min_col, merged.min_col)
        max_col = merged.max_col if max_col is None else max(max_col, merged.max_col)
    if min_row is None:
        return 1, 1, 1, 1
    return min_row, max_row, min_col, max_col


def _normalize_month_token(text: str) -> str:
    cleaned = "".join(ch for ch in text.upper() if ch.isalpha())
    return cleaned


def _find_weight_col(
    ws: Worksheet,
    month_row: int | None,
    min_col: int,
    max_col: int,
) -> int | None:
    """Locate the WEIGHT column that sits just right of the colored category bars."""
    row_start = max(1, (month_row or 5) - 2)
    row_end = (month_row or 5) + 4
    for row in range(row_start, row_end + 1):
        for col in range(min_col, max_col + 1):
            if _normalize_label(_cell_text(ws.cell(row, col))) == "weight":
                return col
    return None


def _find_category_bar_col(
    ws: Worksheet,
    min_col: int,
    max_col: int,
    body_start: int,
    *,
    weight_col: int | None = None,
) -> int | None:
    """First column of the colored category rail (Safety / Finance / etc.)."""
    limit = weight_col if weight_col is not None else min(max_col, min_col + 6)
    for col in range(min_col, limit + 1):
        for row in range(body_start, body_start + 50):
            text = _cell_text(ws.cell(row, col))
            if _is_category_title(text):
                return col
            fill = _fill_rgb(ws.cell(row, col))
            # Strong colored fill left of WEIGHT = category bar (not the pale % rail).
            if (
                fill
                and _luminance(fill) < 170
                and max(fill) - min(fill) > 25
                and (weight_col is None or col < weight_col)
            ):
                return col
    if weight_col is not None and weight_col > min_col:
        return weight_col - 1
    return None


def _scorecard_content_cols(
    ws: Worksheet,
    month_row: int | None,
    min_col: int,
    max_col: int,
    *,
    body_start: int | None = None,
) -> tuple[int, int]:
    """Trim capture to colored category bars through YE (skip duplicate % column on the left)."""
    end = max_col
    if month_row:
        month_cols: list[int] = []
        for col in range(min_col, max_col + 1):
            token = _normalize_month_token(_cell_text(ws.cell(month_row, col)))
            if token and (token in MONTH_HEADER_TOKENS or token[:3] in MONTH_HEADER_TOKENS):
                month_cols.append(col)
        if month_cols:
            end = max(month_cols)

    weight_col = _find_weight_col(ws, month_row, min_col, max_col)
    # Colored category bars sit immediately left of WEIGHT; columns further left
    # are the duplicate percentage rail (and red border) — skip them.
    if weight_col is not None and weight_col > min_col:
        start = weight_col - 1
    else:
        start = (
            _find_category_bar_col(
                ws,
                min_col,
                max_col,
                body_start or ((month_row + 1) if month_row else min_col),
                weight_col=weight_col,
            )
            or min_col
        )

    if end == max_col:
        sample_rows = list(range(1, min(40, int(ws.max_row or 40)) + 1))
        while end > start:
            if any(_cell_text(ws.cell(r, end)) or _fill_rgb(ws.cell(r, end)) for r in sample_rows):
                break
            end -= 1
    return start, end


def _is_month_footer_row(ws: Worksheet, row: int, min_col: int, max_col: int) -> bool:
    """True for the repeating JAN..DEC footer row under Operations (not a real KPI row)."""
    month_hits = 0
    other = 0
    for col in range(min_col, max_col + 1):
        text = _cell_text(ws.cell(row, col))
        if not text:
            continue
        token = _normalize_month_token(text)
        if token in MONTH_HEADER_TOKENS or token[:3] in MONTH_HEADER_TOKENS:
            month_hits += 1
        else:
            other += 1
    return month_hits >= 6 and other == 0


def _is_opportunities_row(ws: Worksheet, row: int, min_col: int, max_col: int) -> bool:
    """True for the gray OPPORTUNITIES footer under Overall Total Score."""
    for col in range(min_col, max_col + 1):
        text = _cell_text(ws.cell(row, col)).casefold()
        if not text:
            continue
        if text.startswith("opportunit") or "opportunities:" in text or text.startswith("count of"):
            return True
    return False


def _trim_capture_end_row(
    ws: Worksheet,
    start_row: int,
    end_row: int,
    min_col: int,
    max_col: int,
) -> int:
    """Drop trailing month-footer / OPPORTUNITIES rows from the capture."""
    end = end_row
    while end > start_row:
        if _is_month_footer_row(ws, end, min_col, max_col) or _is_opportunities_row(
            ws, end, min_col, max_col
        ):
            end -= 1
            continue
        break
    return end


def _find_month_header_row(ws: Worksheet, min_row: int, max_row: int, min_col: int, max_col: int) -> int | None:
    """Locate the JAN..DEC / YTD / YE header row used by the System scorecard."""
    best_row = None
    best_hits = 0
    for row in range(min_row, min(max_row, min_row + 40) + 1):
        hits = 0
        for col in range(min_col, max_col + 1):
            token = _normalize_month_token(_cell_text(ws.cell(row, col)))
            if not token:
                continue
            if token in MONTH_HEADER_TOKENS or token[:3] in MONTH_HEADER_TOKENS:
                hits += 1
        if hits > best_hits:
            best_hits = hits
            best_row = row
    if best_hits >= 6:
        return best_row
    return None


def _header_band(ws: Worksheet, month_row: int, min_col: int, max_col: int, first_section_row: int) -> tuple[int, int]:
    """Include only the month header row (and an immediate label row), not the whole body."""
    start = month_row
    if month_row > 1 and _row_has_content(ws, month_row - 1, min_col, max_col):
        # Keep a short title row above months, but never more than 1 row.
        start = month_row - 1
    end = month_row
    # At most one WEIGHT/KPI label row directly under the months.
    probe = month_row + 1
    if probe < first_section_row:
        texts = " ".join(_cell_text(ws.cell(probe, c)).casefold() for c in range(min_col, min(min_col + 8, max_col + 1)))
        if "weight" in texts or texts.strip() == "kpi" or "total score" in texts:
            end = probe
    return start, end


def _is_category_title(text: str) -> bool:
    if not text or len(text) > 80:
        return False
    lower = text.casefold().strip()
    if lower in METRIC_LABELS or lower in KNOWN_KPI_NAMES:
        return False
    # Avoid matching bare "safety" inside KPI names like "Psych. Safety".
    if lower.startswith("psych"):
        return False
    token = _normalize_month_token(text)
    if token in MONTH_HEADER_TOKENS or token[:3] in MONTH_HEADER_TOKENS:
        return False
    if "(" in text and "%" in text:
        return True
    if any(k in lower for k in CATEGORY_KEYWORDS):
        return True
    return False


def _left_label_columns(min_col: int, max_col: int) -> range:
    # Real category rail is the far-left columns only (not the KPI name column).
    return range(min_col, min(min_col + 2, max_col + 1))


def _assign_default_titles(
    sections: list[tuple[int, int, str, bool]],
) -> list[tuple[int, int, str, bool]]:
    """Fill missing/generic titles using the known System category order."""
    out: list[tuple[int, int, str, bool]] = []
    for idx, (start, end, title, is_black) in enumerate(sections):
        cleaned = title.strip()
        generic = (
            not cleaned
            or cleaned.casefold().startswith("section ")
            or cleaned.casefold() in {"system", "system scorecard", "total score"}
        )
        if generic and idx < len(DEFAULT_SYSTEM_CATEGORIES):
            cleaned = DEFAULT_SYSTEM_CATEGORIES[idx]
        elif generic:
            cleaned = f"Section {idx + 1}"
        if idx == len(sections) - 1 and (
            is_black or "overall" in cleaned.casefold() or "opportunit" in cleaned.casefold()
        ):
            is_black = True
            if "overall" not in cleaned.casefold():
                cleaned = DEFAULT_SYSTEM_CATEGORIES[-1]
        out.append((start, end, cleaned, is_black))
    return out


def _vertical_category_merges(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
) -> list[tuple[int, int, str, bool]]:
    """
    Category sections are tall merged cells in the leftmost columns, e.g.
    'Safety & Security (25.0%)' spanning many rows with a maroon/blue/orange fill.
    Ignores sheet-wide wrapper merges like a single 'System' label.
    """
    body_height = max(max_row - min_row + 1, 1)
    found: list[tuple[int, int, str, bool]] = []
    for merged in ws.merged_cells.ranges:
        height = merged.max_row - merged.min_row + 1
        width = merged.max_col - merged.min_col + 1
        if height < 4 or width > 2:
            continue
        if merged.min_col > min_col + 4:
            continue
        if merged.max_row < min_row or merged.min_row > max_row:
            continue
        # Skip wrappers that cover most of the sheet (seen as one 'System' block).
        if height >= max(20, int(body_height * 0.45)):
            continue
        title = _cell_text(ws.cell(merged.min_row, merged.min_col))
        if not title:
            continue
        if title.strip().casefold() in {"system", "system scorecard"} and height >= 20:
            continue
        rgb = _fill_rgb(ws.cell(merged.min_row, merged.min_col))
        is_black = _is_dark(rgb)
        found.append((merged.min_row, merged.max_row, title, is_black))

    found.sort(key=lambda item: item[0])
    return found


def _find_total_score_rows(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
) -> list[int]:
    """Return 1-based rows that contain a TOTAL SCORE label anywhere in the used width."""
    starts: list[int] = []
    for row in range(min_row, max_row + 1):
        for col in range(min_col, max_col + 1):
            if _is_total_score_label(_cell_text(ws.cell(row, col))):
                starts.append(row)
                break
    return starts


def _total_score_rows_via_excel_com(
    workbook_bytes: bytes,
    sheet_name: str,
    min_row: int,
    max_row: int,
) -> list[int]:
    """Read displayed cell text through Excel — catches labels openpyxl misses."""
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return []

    excel = None
    wb = None
    starts: list[int] = []
    try:
        try:
            import pythoncom

            pythoncom.CoInitialize()
        except Exception:
            pythoncom = None  # type: ignore
        with tempfile.TemporaryDirectory(prefix="mpr_scorecard_detect_") as tmp:
            path = Path(tmp) / "scorecards.xlsx"
            path.write_bytes(workbook_bytes)
            excel = _create_excel_application(visible=False)
            wb = excel.Workbooks.Open(str(path), ReadOnly=True, UpdateLinks=0)
            sheet = None
            for candidate in wb.Worksheets:
                if _norm_sheet_key(candidate.Name) == _norm_sheet_key(sheet_name):
                    sheet = candidate
                    break
            if sheet is None:
                return []
            used = sheet.UsedRange
            values = used.Value
            if values is None:
                return []
            # UsedRange.Value is a tuple-of-tuples for multi-cell ranges.
            if not isinstance(values, tuple):
                values = ((values,),)
            elif values and not isinstance(values[0], tuple):
                values = (values,)
            origin_row = int(used.Row)
            for r_idx, row_vals in enumerate(values):
                row_num = origin_row + r_idx
                if row_num < min_row or row_num > max_row:
                    continue
                if not isinstance(row_vals, tuple):
                    row_vals = (row_vals,)
                for value in row_vals:
                    if value is None:
                        continue
                    if _is_total_score_label(str(value)):
                        starts.append(row_num)
                        break
            if starts:
                logger.info("Excel COM found TOTAL SCORE rows: %s", starts)
            return starts
    except Exception as exc:
        logger.info("Excel COM TOTAL SCORE scan unavailable: %s", exc)
        return []
    finally:
        try:
            if wb is not None:
                wb.Close(SaveChanges=False)
        except Exception:
            pass
        try:
            if excel is not None:
                excel.Quit()
        except Exception:
            pass
        try:
            if pythoncom is not None:
                pythoncom.CoUninitialize()
        except Exception:
            pass


def _sections_from_total_score_starts(
    ws: Worksheet,
    starts: list[int],
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
) -> list[tuple[int, int, str, bool]]:
    if len(starts) < 2:
        return []
    # Dedupe + sort; ignore TOTAL SCORE rows above the body.
    starts = sorted({row for row in starts if row >= min_row})
    if len(starts) < 2:
        return []

    sections: list[tuple[int, int, str, bool]] = []
    for idx, start in enumerate(starts):
        end = starts[idx + 1] - 1 if idx + 1 < len(starts) else max_row
        title = ""
        is_black = False
        for row in range(start, min(start + 10, end + 1)):
            for col in _left_label_columns(min_col, max_col):
                cell = ws.cell(row, col)
                text = _cell_text(cell)
                rgb = _fill_rgb(cell)
                if rgb is not None and _is_dark(rgb):
                    is_black = True
                if text and _is_category_title(text):
                    title = text
                    break
            if title:
                break
        if not title and idx < len(DEFAULT_SYSTEM_CATEGORIES):
            title = DEFAULT_SYSTEM_CATEGORIES[idx]
        elif not title:
            title = f"Section {idx + 1}"
        sections.append((start, end, title, is_black))
    return _assign_default_titles(sections)


def _total_score_section_starts(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
    *,
    workbook_bytes: bytes | None = None,
    sheet_name: str = "System",
) -> list[tuple[int, int, str, bool]]:
    """
    Each category block in the live System tab starts with a grey TOTAL SCORE row.
    Use those rows as section boundaries.
    """
    starts = _find_total_score_rows(ws, min_row, max_row, min_col, max_col)
    if len(starts) < 2 and workbook_bytes:
        starts = _total_score_rows_via_excel_com(
            workbook_bytes, sheet_name, min_row, max_row
        )
    return _sections_from_total_score_starts(ws, starts, min_row, max_row, min_col, max_col)


def _default_category_bands(
    min_row: int,
    max_row: int,
) -> list[tuple[int, int, str, bool]]:
    """Last-resort split into the six known System categories by row bands."""
    names = DEFAULT_SYSTEM_CATEGORIES
    height = max_row - min_row + 1
    if height < len(names) * 4:
        return []
    sections: list[tuple[int, int, str, bool]] = []
    for idx, name in enumerate(names):
        start = min_row + (height * idx) // len(names)
        end = min_row + (height * (idx + 1)) // len(names) - 1
        if idx == len(names) - 1:
            end = max_row
        sections.append((start, max(start, end), name, idx == len(names) - 1))
    return sections


def _category_title_rows(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
) -> list[tuple[int, int, str, bool]]:
    """Find category labels by text in the far-left rail only."""
    found: list[tuple[int, int, str, bool]] = []
    seen_rows: set[int] = set()
    for row in range(min_row, max_row + 1):
        for col in _left_label_columns(min_col, max_col):
            text = _cell_text(ws.cell(row, col))
            if not _is_category_title(text):
                continue
            if text.strip().casefold() in {"system", "system scorecard"}:
                continue
            if row in seen_rows:
                continue
            rgb = _fill_rgb(ws.cell(row, col))
            found.append((row, row, text, _is_dark(rgb)))
            seen_rows.add(row)
            break
    return _assign_default_titles(found) if len(found) >= 3 else []


def _fill_run_sections(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
) -> list[tuple[int, int, str, bool]]:
    """Split on left-column fill-color changes across columns A-E."""
    best_col = min_col
    best_filled = -1
    for col in range(min_col, min(min_col + 5, max_col + 1)):
        filled = sum(1 for row in range(min_row, max_row + 1) if _fill_rgb(ws.cell(row, col)) is not None)
        if filled > best_filled:
            best_filled = filled
            best_col = col
    if best_filled < 8:
        return []

    sections: list[tuple[int, int, str, bool]] = []
    current_start = None
    current_rgb = None
    current_title = ""
    current_black = False

    def close(end_row: int) -> None:
        nonlocal current_start, current_rgb, current_title, current_black
        if current_start is None:
            return
        sections.append(
            (
                current_start,
                end_row,
                current_title or f"Section {len(sections) + 1}",
                current_black,
            )
        )
        current_start = None
        current_rgb = None
        current_title = ""
        current_black = False

    for row in range(min_row, max_row + 1):
        rgb = _fill_rgb(ws.cell(row, best_col))
        text = _cell_text(ws.cell(row, best_col))
        if rgb is None:
            continue
        if current_start is None:
            current_start = row
            current_rgb = rgb
            current_title = text if _is_category_title(text) else ""
            current_black = _is_dark(rgb)
            continue
        if rgb != current_rgb:
            close(row - 1)
            current_start = row
            current_rgb = rgb
            current_title = text if _is_category_title(text) else ""
            current_black = _is_dark(rgb)
        elif not current_title and _is_category_title(text):
            current_title = text
            current_black = current_black or _is_dark(rgb)

    if current_start is not None:
        close(max_row)

    kept = [s for s in sections if s[1] - s[0] >= 3]
    return _assign_default_titles(kept) if len(kept) >= 3 else []


def _choose_raw_sections(
    ws: Worksheet,
    min_row: int,
    max_row: int,
    min_col: int,
    max_col: int,
    *,
    workbook_bytes: bytes | None = None,
    sheet_name: str = "System",
) -> list[tuple[int, int, str, bool]]:
    """Try several strategies; prefer TOTAL SCORE boundaries for the live workbook."""

    def total_score() -> list[tuple[int, int, str, bool]]:
        return _total_score_section_starts(
            ws,
            min_row,
            max_row,
            min_col,
            max_col,
            workbook_bytes=workbook_bytes,
            sheet_name=sheet_name,
        )

    candidates = [
        ("TOTAL SCORE rows", total_score),
        ("vertical merges", lambda: _vertical_category_merges(ws, min_row, max_row, min_col, max_col)),
        ("left fill runs", lambda: _fill_run_sections(ws, min_row, max_row, min_col, max_col)),
        ("category titles", lambda: _category_title_rows(ws, min_row, max_row, min_col, max_col)),
        ("default 6-band split", lambda: _default_category_bands(min_row, max_row)),
    ]
    for name, factory in candidates:
        found = factory()
        # Prefer a full System scorecard (5–6 categories). Accept 3+ as usable.
        if len(found) >= 3:
            if name == "default 6-band split":
                logger.warning(
                    "Using even 6-category row split R%s-%s (TOTAL SCORE labels not readable)",
                    min_row,
                    max_row,
                )
            logger.info("System section strategy=%s count=%s", name, len(found))
            return _assign_default_titles(found)
        logger.info("System section strategy=%s count=%s (skip)", name, len(found))
    return []


def detect_system_layout(
    workbook_bytes: bytes,
    sheet_name: str = "System",
) -> SystemLayout:
    """
    Parse the System scorecard layout matching the live GSE file:
    month header (JAN..YE) + vertical category bars on the left.
    """
    # data_only=False keeps labels/formulas visible for section detection.
    wb = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
    if sheet_name not in wb.sheetnames:
        match = _find_sheet_name(list(wb.sheetnames), sheet_name)
        if match is None:
            raise ValueError(f"Sheet {sheet_name!r} not found. Available: {wb.sheetnames}")
        sheet_name = match
    ws = wb[sheet_name]
    min_row, max_row, min_col, max_col = _used_bounds(ws)

    month_row = _find_month_header_row(ws, min_row, max_row, min_col, max_col)
    body_start = (month_row + 1) if month_row else min_row

    raw_sections = _choose_raw_sections(
        ws,
        body_start,
        max_row,
        min_col,
        max_col,
        workbook_bytes=workbook_bytes,
        sheet_name=sheet_name,
    )
    if not raw_sections:
        logger.warning(
            "Could not split System categories (cols %s-%s rows %s-%s); capturing full body",
            get_column_letter(min_col),
            get_column_letter(max_col),
            body_start,
            max_row,
        )
        raw_sections = [(body_start, max_row, "System", False)]

    # Extend each section to the row before the next category starts.
    normalized: list[tuple[int, int, str, bool]] = []
    for idx, (start, end, title, is_black) in enumerate(raw_sections):
        next_start = raw_sections[idx + 1][0] if idx + 1 < len(raw_sections) else max_row + 1
        end = max(end, next_start - 1)
        while end > start and not _row_has_content(ws, end, min_col, max_col):
            end -= 1
        # Drop duplicate month footer / OPPORTUNITIES footer rows.
        while end > start and (
            _is_month_footer_row(ws, end, min_col, max_col)
            or _is_opportunities_row(ws, end, min_col, max_col)
        ):
            end -= 1
        # Mark trailing dark totals block as black when fill says so.
        if not is_black:
            for col in _left_label_columns(min_col, max_col):
                if _is_dark(_fill_rgb(ws.cell(start, col))):
                    is_black = True
                    break
        normalized.append((start, end, title, is_black))

    first_section_row = normalized[0][0]
    if month_row:
        header_start, header_end = _header_band(ws, month_row, min_col, max_col, first_section_row)
    else:
        header_start = min_row
        header_end = max(min_row, first_section_row - 1)

    # Start at colored category bars (skip duplicate % rail); end at YE.
    content_min, content_max = _scorecard_content_cols(
        ws,
        month_row,
        min_col,
        max_col,
        body_start=first_section_row,
    )

    sections = [
        Section(
            index=idx,
            title=title,
            start_row=start,
            end_row=end,
            start_col=content_min,
            end_col=content_max,
            is_black=is_black,
        )
        for idx, (start, end, title, is_black) in enumerate(normalized)
    ]

    layout = SystemLayout(
        header_start_row=header_start,
        header_end_row=header_end,
        start_col=content_min,
        end_col=content_max,
        sections=sections,
    )
    wb.close()
    logger.info(
        "System layout: header rows %s-%s cols %s-%s; sections=%s",
        layout.header_start_row,
        layout.header_end_row,
        get_column_letter(layout.start_col),
        get_column_letter(layout.end_col),
        [(s.index, s.title, f"R{s.start_row}-{s.end_row}", s.is_black) for s in sections],
    )
    return layout


def detect_system_sections(
    workbook_bytes: bytes,
    sheet_name: str = "System",
) -> list[Section]:
    """Backward-compatible helper returning only category sections."""
    return detect_system_layout(workbook_bytes, sheet_name=sheet_name).sections


def select_sections_for_slide(
    sections: list[Section],
    *,
    mode: str,
    count: int = 3,
    include_black: bool = False,
    match: list[str] | None = None,
    indices: list[int] | None = None,
) -> list[Section]:
    """Choose which sections belong on a System Scorecard slide."""
    if not sections:
        return []

    if mode == "indices":
        chosen = [sections[i] for i in (indices or []) if 0 <= i < len(sections)]
        return sorted({s.index: s for s in chosen}.values(), key=lambda s: s.start_row)

    if mode == "match":
        patterns = [p.strip() for p in (match or []) if str(p).strip()]
        if not patterns:
            raise ValueError("sections: match requires a non-empty match list")
        chosen: list[Section] = []
        used: set[int] = set()
        for pattern in patterns:
            needle = pattern.casefold()
            for section in sections:
                if section.index in used:
                    continue
                title = section.title.casefold()
                if needle in title or title in needle:
                    chosen.append(section)
                    used.add(section.index)
                    break
        chosen = sorted({s.index: s for s in chosen}.values(), key=lambda s: s.start_row)
        if chosen:
            return chosen
        if indices:
            logger.warning(
                "Name match failed for %s; falling back to indices %s (available=%s)",
                patterns,
                indices,
                [s.title for s in sections],
            )
            fallback = select_sections_for_slide(sections, mode="indices", indices=indices)
            if fallback:
                return fallback
            if len(sections) >= 6 and indices:
                # Canonical System layout: 0-2 slide3, 3-5 slide4.
                start = 0 if min(indices) < 3 else 3
                return sections[start : start + 3]
            if len(sections) == 1:
                return list(sections)
            mid = max(1, len(sections) // 2)
            if min(indices) >= mid:
                return sections[mid:]
            return sections[:mid]
        logger.warning(
            "No System sections matched %s. Available: %s",
            patterns,
            [s.title for s in sections],
        )
        return chosen

    black = [s for s in sections if s.is_black]
    non_black = [s for s in sections if not s.is_black] or list(sections)

    if mode == "first":
        chosen = non_black[:count]
    elif mode == "last":
        chosen = non_black[-count:] if len(non_black) >= count else list(non_black)
        if include_black:
            for b in black:
                if b not in chosen:
                    chosen.append(b)
            if not black and sections[-1] not in chosen:
                chosen.append(sections[-1])
    else:
        raise ValueError(f"Unknown section mode: {mode!r}")

    chosen = sorted(chosen, key=lambda s: s.start_row)
    return chosen


_LAYOUT_CACHE: dict[tuple[int, str], SystemLayout] = {}


def _cached_system_layout(workbook_bytes: bytes, sheet_name: str) -> SystemLayout:
    key = (hash(workbook_bytes), _norm_sheet_key(sheet_name))
    cached = _LAYOUT_CACHE.get(key)
    if cached is not None:
        return cached
    layout = detect_system_layout(workbook_bytes, sheet_name=sheet_name)
    _LAYOUT_CACHE[key] = layout
    return layout


def _range_address(start_row: int, end_row: int, start_col: int, end_col: int) -> str:
    return f"{get_column_letter(start_col)}{start_row}:{get_column_letter(end_col)}{end_row}"


def _png_from_pil(img) -> bytes:
    if img.mode == "RGBA":
        background = Image.new("RGB", img.size, (255, 255, 255))
        background.paste(img, mask=img.split()[3])
        img = background
    elif img.mode != "RGB":
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="PNG", optimize=False)
    return buf.getvalue()


def upscale_png_to_min_width(
    png_bytes: bytes,
    *,
    min_width: int = 1200,
    max_scale: float = 4.0,
) -> bytes:
    """LANCZOS-upscale a small capture so legend text stays sharp on the slide.

    Excel CopyPicture of tiny legend tables often returns ~400–600px-wide bitmaps.
    Placing those into even a modest legend slot looks pixelated; pre-upscaling
    gives PowerPoint more pixels to work with.
    """
    if min_width <= 0:
        return png_bytes
    with Image.open(io.BytesIO(png_bytes)) as img:
        img = img.convert("RGB")
        if img.width >= min_width:
            return png_bytes
        scale = min(max_scale, min_width / max(img.width, 1))
        if scale <= 1.01:
            return png_bytes
        new_size = (
            max(1, int(img.width * scale)),
            max(1, int(img.height * scale)),
        )
        upscaled = img.resize(new_size, Image.Resampling.LANCZOS)
        logger.info(
            "Upscaled capture %sx%s -> %sx%s (min_width=%s)",
            img.width,
            img.height,
            upscaled.width,
            upscaled.height,
            min_width,
        )
        return _png_from_pil(upscaled)


def _clipboard_image():
    """Read an image from the Windows clipboard, if present."""
    from PIL import ImageGrab

    return ImageGrab.grabclipboard()


def _is_mostly_blank(png_bytes: bytes, *, white_ratio: float = 0.90) -> bool:
    """True when a capture is almost entirely white (the failure mode on slides 3/4)."""
    with Image.open(io.BytesIO(png_bytes)) as img:
        img = img.convert("RGB")
        sample = img.resize((max(1, img.width // 10), max(1, img.height // 10)))
        pixels = list(sample.getdata())
    if not pixels:
        return True
    whiteish = sum(1 for r, g, b in pixels if r >= 245 and g >= 245 and b >= 245)
    return (whiteish / len(pixels)) >= white_ratio


def _autocrop_whitespace(png_bytes: bytes, *, threshold: int = 248, pad: int = 4) -> bytes:
    """Trim near-white margins so the scorecard fills the picture like the template EMF."""
    with Image.open(io.BytesIO(png_bytes)) as img:
        rgb = img.convert("RGB")
        w, h = rgb.size
        pixels = rgb.load()

        def is_content(x: int, y: int) -> bool:
            r, g, b = pixels[x, y]
            return r < threshold or g < threshold or b < threshold

        top = 0
        while top < h and not any(is_content(x, top) for x in range(w)):
            top += 1
        bottom = h - 1
        while bottom > top and not any(is_content(x, bottom) for x in range(w)):
            bottom -= 1
        left = 0
        while left < w and not any(is_content(left, y) for y in range(top, bottom + 1)):
            left += 1
        right = w - 1
        while right > left and not any(is_content(right, y) for y in range(top, bottom + 1)):
            right -= 1

        if right - left < 20 or bottom - top < 20:
            return png_bytes
        box = (
            max(0, left - pad),
            max(0, top - pad),
            min(w, right + pad + 1),
            min(h, bottom + pad + 1),
        )
        return _png_from_pil(rgb.crop(box))


def _validate_capture(
    png_bytes: bytes,
    *,
    label: str,
    min_w: int = 400,
    min_h: int = 150,
    require_wide: bool = True,
) -> bytes:
    png_bytes = _autocrop_whitespace(png_bytes)
    with Image.open(io.BytesIO(png_bytes)) as img:
        width, height = img.width, img.height
        png = _png_from_pil(img)
    if width < min_w or height < min_h:
        raise RuntimeError(f"{label} capture too small: {width}x{height}")
    if _is_mostly_blank(png):
        raise RuntimeError(f"{label} capture is mostly blank/white ({width}x{height}, {len(png)} bytes)")
    # Template System EMFs are wide scorecards — optional for smaller GIR panels.
    aspect = width / max(height, 1)
    if require_wide and aspect < 1.2:
        raise RuntimeError(
            f"{label} capture aspect {aspect:.2f} looks wrong for a scorecard "
            f"({width}x{height}) — expected a wide grid like the template"
        )
    logger.info("%s capture OK: %sx%s aspect=%.2f (%s bytes)", label, width, height, aspect, len(png))
    return png


def _copy_picture_to_png(rng, *, appearance: int, fmt: int) -> bytes | None:
    try:
        rng.CopyPicture(Appearance=appearance, Format=fmt)
    except Exception:
        return None
    for _ in range(40):
        time.sleep(0.15)
        grabbed = _clipboard_image()
        if grabbed is not None:
            return _png_from_pil(grabbed)
    return None


def _clear_broken_win32com_cache() -> None:
    """Remove corrupted pywin32 gen_py cache (fixes CLSIDToPackageMap errors)."""
    try:
        import win32com  # type: ignore

        gen_path = Path(getattr(win32com, "__gen_path__", "") or "")
        if gen_path.is_dir():
            import shutil

            shutil.rmtree(gen_path, ignore_errors=True)
            logger.warning("Cleared broken win32com gen_py cache at %s", gen_path)
            print(f">>> Cleared broken Excel COM cache: {gen_path}")
    except Exception as exc:
        logger.warning("Could not clear win32com cache: %s", exc)

    # Drop cached modules so the next Dispatch rebuilds cleanly.
    try:
        import sys

        doomed = [name for name in list(sys.modules) if name.startswith("win32com.gen_py")]
        for name in doomed:
            sys.modules.pop(name, None)
    except Exception:
        pass


def _create_excel_application(*, visible: bool = True):
    """Create Excel.Application via late-bound COM (avoids broken gencache).

    The classic failure mode is:
      module 'win32com.gen_py....' has no attribute 'CLSIDToPackageMap'
    which makes EnsureDispatch unusable and forces Pillow fallbacks. Prefer
    DispatchEx / dynamic.Dispatch, and rebuild the cache if needed.
    """
    import win32com.client  # type: ignore

    try:
        import pythoncom

        pythoncom.CoInitialize()
    except Exception:
        pass

    errors: list[str] = []

    def _configure(excel):
        excel.Visible = bool(visible)
        excel.DisplayAlerts = False
        excel.ScreenUpdating = True
        try:
            excel.AskToUpdateLinks = False
        except Exception:
            pass
        if visible:
            try:
                # Normal window — maximized+window-grab looks like Excel UI.
                excel.WindowState = -4143  # xlNormal
            except Exception:
                pass
            try:
                excel.Width = 1400
                excel.Height = 900
            except Exception:
                pass
        return excel

    # 1) Late-bound DispatchEx — does not depend on gen_py early binding.
    for factory_name, factory in (
        ("DispatchEx", win32com.client.DispatchEx),
        ("dynamic.Dispatch", win32com.client.dynamic.Dispatch),
        ("Dispatch", win32com.client.Dispatch),
    ):
        try:
            excel = factory("Excel.Application")
            return _configure(excel)
        except Exception as exc:
            errors.append(f"{factory_name}:{exc}")
            if "CLSIDToPackageMap" in str(exc):
                _clear_broken_win32com_cache()

    # 2) Last resort: wipe cache and try DispatchEx once more.
    _clear_broken_win32com_cache()
    try:
        excel = win32com.client.DispatchEx("Excel.Application")
        return _configure(excel)
    except Exception as exc:
        errors.append(f"retry DispatchEx:{exc}")
        raise RuntimeError(
            "Could not start Excel via win32com. "
            "Try: pip install --force-reinstall pywin32 && "
            r"Remove-Item -Recurse -Force $env:LOCALAPPDATA\Temp\gen_py "
            f"(errors: {'; '.join(errors)})"
        ) from exc


def _open_excel_workbook(workbook_path: Path):
    """Open Excel visible (required for non-blank CopyPicture) and return (excel, wb)."""
    last_exc: Exception | None = None
    for attempt in range(2):
        excel = None
        try:
            excel = _create_excel_application(visible=True)
            wb = excel.Workbooks.Open(
                str(workbook_path.resolve()),
                ReadOnly=True,
                UpdateLinks=0,
            )
            if attempt > 0:
                print(">>> Excel COM reopen succeeded after cache clear")
            return excel, wb
        except Exception as exc:
            last_exc = exc
            try:
                if excel is not None:
                    excel.Quit()
            except Exception:
                pass
            if attempt == 0 and "CLSIDToPackageMap" in str(exc):
                logger.warning("Excel open hit corrupt COM cache; clearing and retrying")
                print(">>> Excel COM cache corrupt on open — clearing gen_py and retrying...")
                _clear_broken_win32com_cache()
                continue
            raise
    raise RuntimeError(f"Could not open Excel workbook: {last_exc}")




def _prepare_sheet_for_copy(excel, ws, range_addr: str, *, zoom: int = 100):
    """Select range at the given zoom (100% matches template scorecard EMFs).

    Use zoom > 100 for small ranges (e.g. legends) so CopyPicture returns denser pixels.
    """
    ws.Activate()
    try:
        excel.ActiveWindow.DisplayGridlines = False
        excel.ActiveWindow.DisplayHeadings = False
        excel.ActiveWindow.Zoom = int(max(50, min(400, zoom)))
    except Exception:
        pass
    try:
        excel.CalculateUntilAsyncQueriesDone()
    except Exception:
        try:
            excel.Calculate()
        except Exception:
            pass

    rng = ws.Range(range_addr)
    try:
        excel.Goto(rng, True)
    except Exception:
        try:
            rng.Select()
        except Exception:
            pass
    time.sleep(0.35)
    return rng


def _copy_range_picture(rng) -> bytes:
    """CopyPicture the Excel range to a validated PNG (template-style paste)."""
    errors: list[str] = []
    # Prefer printer quality, then screen bitmap. Avoid window screenshots.
    for appearance, fmt in ((2, 2), (1, 2)):  # xlPrinter/xlScreen + xlBitmap
        try:
            png = _copy_picture_to_png(rng, appearance=appearance, fmt=fmt)
            if png:
                return _validate_capture(png, label=f"CopyPicture({appearance}/{fmt})")
        except Exception as exc:
            errors.append(f"{appearance}/{fmt}:{exc}")

    # Chart.Export fallback still uses the range picture, not the Excel window.
    chart_obj = None
    export_path = Path(tempfile.gettempdir()) / f"mpr_scorecard_{os.getpid()}_{time.time_ns()}.png"
    try:
        ws = rng.Worksheet
        width = max(float(rng.Width), 40.0)
        height = max(float(rng.Height), 40.0)
        max_side = 1600.0
        scale = min(1.0, max_side / width, max_side / height)
        chart_obj = ws.ChartObjects().Add(10, 10, width * scale, height * scale)
        for appearance, fmt in ((2, 2), (1, 2)):
            try:
                rng.CopyPicture(Appearance=appearance, Format=fmt)
                time.sleep(0.25)
                chart_obj.Chart.Paste()
                time.sleep(0.25)
                if export_path.exists():
                    export_path.unlink()
                chart_obj.Chart.Export(str(export_path))
                if export_path.exists() and export_path.stat().st_size > 1500:
                    return _validate_capture(
                        export_path.read_bytes(),
                        label=f"Chart.Export({appearance}/{fmt})",
                    )
            except Exception as exc:
                errors.append(f"chart:{exc}")
    finally:
        if chart_obj is not None:
            try:
                chart_obj.Delete()
            except Exception:
                pass
        try:
            if export_path.exists():
                export_path.unlink()
        except Exception:
            pass

    raise RuntimeError(f"CopyPicture failed ({'; '.join(errors) or 'no image'})")


def _find_com_worksheet(wb, sheet_name: str):
    for sheet in wb.Worksheets:
        if _norm_sheet_key(sheet.Name) == _norm_sheet_key(sheet_name):
            return sheet
    raise ValueError(f"Worksheet {sheet_name!r} not found in Excel COM open")


def _capture_via_excel_com(
    workbook_path: Path,
    sheet_name: str,
    range_addr: str,
    *,
    zoom: int = 100,
) -> bytes:
    """True Excel range screenshot via CopyPicture (same approach as the template EMFs)."""
    try:
        import win32com.client  # noqa: F401
    except ImportError as exc:
        raise RuntimeError(
            "pywin32 is not installed. Run: pip install pywin32\n"
            "Then re-run main.py for template-quality System scorecard screenshots."
        ) from exc

    last_exc: Exception | None = None
    for attempt in range(2):
        excel = None
        wb = None
        try:
            excel, wb = _open_excel_workbook(workbook_path)
            ws = _find_com_worksheet(wb, sheet_name)
            rng = _prepare_sheet_for_copy(excel, ws, range_addr, zoom=zoom)
            return _copy_range_picture(rng)
        except Exception as exc:
            last_exc = exc
            if attempt == 0 and "CLSIDToPackageMap" in str(exc):
                logger.warning("Excel COM cache corrupt; clearing and retrying once")
                print(">>> Excel COM cache corrupt — clearing gen_py and retrying...")
                _clear_broken_win32com_cache()
                continue
            raise
        finally:
            try:
                if wb is not None:
                    wb.Close(SaveChanges=False)
            except Exception:
                pass
            try:
                if excel is not None:
                    excel.Quit()
            except Exception:
                pass
    raise RuntimeError(f"Excel COM capture failed: {last_exc}")


def _capture_system_sections_via_com(
    workbook_path: Path,
    sheet_name: str,
    *,
    header_start: int,
    header_end: int,
    body_start: int,
    body_end: int,
    start_col: int,
    end_col: int,
) -> bytes:
    """Capture header + selected System sections as one clean scorecard image.

    When People/Finance sit below Safety/CX/Ops, hide the intervening rows so
    CopyPicture returns header+selected blocks contiguous — matching the template
    EMF layout instead of a wrong Pillow crop or an Excel-window screenshot.
    """
    try:
        import win32com.client  # noqa: F401
    except ImportError as exc:
        raise RuntimeError("pywin32 is not installed") from exc

    excel = None
    wb = None
    hidden_rows: list[tuple[int, int]] = []
    try:
        excel, wb = _open_excel_workbook(workbook_path)
        ws = _find_com_worksheet(wb, sheet_name)

        # Hide columns left of the colored category rail (duplicate % column + red border).
        if start_col > 1:
            try:
                ws.Range(
                    ws.Cells(1, 1),
                    ws.Cells(1, start_col - 1),
                ).EntireColumn.Hidden = True
                logger.info(
                    "Hid System columns %s-%s (duplicate %% rail / left border)",
                    get_column_letter(1),
                    get_column_letter(start_col - 1),
                )
            except Exception as exc:
                logger.warning("Could not hide left System columns: %s", exc)

        # Hide rows between the month header and the first requested section.
        gap_start = header_end + 1
        gap_end = body_start - 1
        if gap_end >= gap_start:
            try:
                hide_rng = ws.Range(f"{gap_start}:{gap_end}")
                hide_rng.EntireRow.Hidden = True
                hidden_rows.append((gap_start, gap_end))
                logger.info(
                    "Hid System rows %s-%s so CopyPicture matches template (header + selected sections)",
                    gap_start,
                    gap_end,
                )
            except Exception as exc:
                logger.warning("Could not hide intervening System rows: %s", exc)

        range_addr = _range_address(header_start, body_end, start_col, end_col)
        rng = _prepare_sheet_for_copy(excel, ws, range_addr)
        png = _copy_range_picture(rng)
        logger.info(
            "System COM section capture %s!%s (header R%s-%s + body R%s-%s, cols %s-%s)",
            sheet_name,
            range_addr,
            header_start,
            header_end,
            body_start,
            body_end,
            get_column_letter(start_col),
            get_column_letter(end_col),
        )
        return png
    finally:
        for start, end in hidden_rows:
            try:
                if wb is not None:
                    ws = _find_com_worksheet(wb, sheet_name)
                    ws.Range(f"{start}:{end}").EntireRow.Hidden = False
            except Exception:
                pass
        try:
            if wb is not None:
                wb.Close(SaveChanges=False)
        except Exception:
            pass
        try:
            if excel is not None:
                excel.Quit()
        except Exception:
            pass


def _font(size: int = 12):
    for name in (
        "C:/Windows/Fonts/arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    ):
        path = Path(name)
        if path.exists():
            try:
                return ImageFont.truetype(str(path), size=size)
            except OSError:
                continue
    return ImageFont.load_default()


def _capture_via_render(
    workbook_bytes: bytes,
    sheet_name: str,
    start_row: int,
    end_row: int,
    start_col: int,
    end_col: int,
    *,
    render_scale: int | None = None,
) -> bytes:
    """Approximate Excel look when Windows Excel COM is unavailable or blank."""
    wb_vals = load_workbook(io.BytesIO(workbook_bytes), data_only=True)
    wb_styles = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
    match = _find_sheet_name(list(wb_vals.sheetnames), sheet_name) or sheet_name
    ws_vals = wb_vals[match]
    ws_styles = wb_styles[match]

    col_widths = []
    for c in range(start_col, end_col + 1):
        letter = get_column_letter(c)
        width = ws_styles.column_dimensions[letter].width
        col_widths.append(int(max(float(width or 10.0), 4.0) * 10))
    row_heights = []
    for r in range(start_row, end_row + 1):
        height = ws_styles.row_dimensions[r].height
        row_heights.append(int(max(float(height or 15.0), 12.0) * 1.5))

    img_w = max(sum(col_widths), 40)
    img_h = max(sum(row_heights), 40)
    # Upscale small tables so they remain readable when stretched onto the slide.
    # Tiny legend ranges need a higher multiplier than full scorecards.
    if render_scale is not None:
        scale = max(1, int(render_scale))
    elif img_w < 350 or img_h < 180:
        scale = 4
    elif img_w < 900 or img_h < 500:
        scale = 2
    else:
        scale = 1
    image = Image.new("RGB", (img_w * scale, img_h * scale), (255, 255, 255))
    draw = ImageDraw.Draw(image)
    font = _font(11 * scale)
    bold = _font(12 * scale)

    y = 0
    for r_idx, r in enumerate(range(start_row, end_row + 1)):
        x = 0
        rh = row_heights[r_idx] * scale
        for c_idx, c in enumerate(range(start_col, end_col + 1)):
            cw = col_widths[c_idx] * scale
            style_cell = ws_styles.cell(r, c)
            value_cell = ws_vals.cell(r, c)
            fill = _fill_rgb(style_cell) or (255, 255, 255)
            draw.rectangle([x, y, x + cw, y + rh], fill=fill, outline=(180, 180, 180))
            text = _cell_text(value_cell) or _cell_text(style_cell)
            # Skip raw formula text — it clutters the scorecard render.
            if text.startswith("="):
                text = ""
            if text:
                text_rgb = (
                    _color_rgb(style_cell.font.color)
                    if style_cell.font and style_cell.font.color
                    else None
                )
                if text_rgb is None:
                    text_rgb = (255, 255, 255) if _is_dark(fill) else (20, 20, 20)
                use_font = bold if (style_cell.font and style_cell.font.bold) else font
                draw.text(
                    (x + 3 * scale, y + max(2, (rh - 12 * scale) // 2)),
                    text[:50],
                    fill=text_rgb,
                    font=use_font,
                )
            x += cw
        y += rh

    wb_vals.close()
    wb_styles.close()
    return _png_from_pil(image)


def capture_range_png(
    workbook_bytes: bytes,
    sheet_name: str,
    start_row: int,
    end_row: int,
    start_col: int,
    end_col: int,
    *,
    prefer_excel_com: bool = True,
    zoom: int = 100,
    render_scale: int | None = None,
    min_width: int | None = None,
) -> bytes:
    """Screenshot (or render) an Excel A1-style range to PNG bytes.

    zoom: Excel ActiveWindow zoom for COM CopyPicture (use 150–200 for small legends).
    render_scale: Pillow fallback pixel multiplier (auto when None).
    min_width: optional LANCZOS upscale floor so tiny captures stay readable on slide.
    """
    range_addr = _range_address(start_row, end_row, start_col, end_col)
    if prefer_excel_com:
        try:
            with tempfile.TemporaryDirectory(prefix="mpr_scorecard_") as tmp:
                path = Path(tmp) / "scorecards.xlsx"
                path.write_bytes(workbook_bytes)
                png = _capture_via_excel_com(path, sheet_name, range_addr, zoom=zoom)
                # Extra guard: never place a blank white bitmap on the deck.
                png = _validate_capture(png, label=f"Excel COM {sheet_name}!{range_addr}")
                if min_width:
                    png = upscale_png_to_min_width(png, min_width=min_width)
                with Image.open(io.BytesIO(png)) as img:
                    logger.info(
                        "Captured %s!%s via Excel COM (%s bytes, %sx%s, zoom=%s)",
                        sheet_name,
                        range_addr,
                        len(png),
                        img.width,
                        img.height,
                        zoom,
                    )
                return png
        except Exception as exc:
            msg = str(exc)
            if "CLSIDToPackageMap" in msg:
                logger.warning(
                    "Excel COM cache corrupt (%s); clearing gen_py and retrying once",
                    exc,
                )
                print(">>> Excel COM cache corrupt — clearing gen_py and retrying...")
                _clear_broken_win32com_cache()
                try:
                    with tempfile.TemporaryDirectory(prefix="mpr_scorecard_") as tmp:
                        path = Path(tmp) / "scorecards.xlsx"
                        path.write_bytes(workbook_bytes)
                        png = _capture_via_excel_com(
                            path, sheet_name, range_addr, zoom=zoom
                        )
                        png = _validate_capture(
                            png, label=f"Excel COM retry {sheet_name}!{range_addr}"
                        )
                        if min_width:
                            png = upscale_png_to_min_width(png, min_width=min_width)
                        with Image.open(io.BytesIO(png)) as img:
                            logger.info(
                                "Captured %s!%s via Excel COM retry (%s bytes, %sx%s, zoom=%s)",
                                sheet_name,
                                range_addr,
                                len(png),
                                img.width,
                                img.height,
                                zoom,
                            )
                        return png
                except Exception as retry_exc:
                    logger.warning(
                        "Excel COM retry also failed (%s). Using Pillow fallback.",
                        retry_exc,
                    )
            else:
                logger.warning(
                    "Excel COM capture unavailable or blank (%s). "
                    "Using Pillow cell render fallback so the slide is not empty.",
                    exc,
                )

    png = _capture_via_render(
        workbook_bytes,
        sheet_name,
        start_row,
        end_row,
        start_col,
        end_col,
        render_scale=render_scale,
    )
    if min_width:
        png = upscale_png_to_min_width(png, min_width=min_width)
    with Image.open(io.BytesIO(png)) as img:
        logger.info(
            "Rendered %s!%s via Pillow (%s bytes, %sx%s)",
            sheet_name,
            range_addr,
            len(png),
            img.width,
            img.height,
        )
    if _is_mostly_blank(png):
        logger.error(
            "Pillow render for %s!%s is also mostly blank — check the System sheet has values",
            sheet_name,
            range_addr,
        )
    return png


def _stitch_vertical_many(png_parts: list[bytes]) -> bytes:
    images = [Image.open(io.BytesIO(p)).convert("RGB") for p in png_parts if p]
    if not images:
        raise ValueError("No images to stitch")
    width = max(im.width for im in images)
    height = sum(im.height for im in images)
    canvas = Image.new("RGB", (width, height), (255, 255, 255))
    y = 0
    for im in images:
        canvas.paste(im, (0, y))
        y += im.height
    buf = io.BytesIO()
    canvas.save(buf, format="PNG")
    return buf.getvalue()


def _sections_are_contiguous(sections: list[Section]) -> bool:
    ordered = sorted(sections, key=lambda s: s.start_row)
    for idx in range(1, len(ordered)):
        if ordered[idx].start_row > ordered[idx - 1].end_row + 2:
            return False
    return True


def _row_height_weights(
    workbook_bytes: bytes,
    sheet_name: str,
    start_row: int,
    end_row: int,
) -> list[float]:
    """Relative Excel row heights used to crop bands from a full-range screenshot."""
    wb = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
    try:
        match = _find_sheet_name(list(wb.sheetnames), sheet_name)
        if match is None:
            raise ValueError(f"Sheet {sheet_name!r} not found. Available: {wb.sheetnames}")
        ws = wb[match]
        heights: list[float] = []
        for row in range(start_row, end_row + 1):
            height = ws.row_dimensions[row].height
            heights.append(max(float(height or 15.0), 8.0))
        return heights
    finally:
        wb.close()


def _crop_row_band(
    png_bytes: bytes,
    heights: list[float],
    start_idx: int,
    end_idx: int,
) -> bytes:
    """Crop an inclusive 0-based row-index band from a PNG of those row heights."""
    if not heights or start_idx > end_idx:
        raise ValueError("Invalid crop band")
    start_idx = max(0, start_idx)
    end_idx = min(len(heights) - 1, end_idx)
    total = sum(heights)
    if total <= 0:
        return png_bytes
    with Image.open(io.BytesIO(png_bytes)) as image:
        image = image.convert("RGB")
        y0 = int(round(image.height * sum(heights[:start_idx]) / total))
        y1 = int(round(image.height * sum(heights[: end_idx + 1]) / total))
        y0 = max(0, min(image.height - 1, y0))
        y1 = max(y0 + 1, min(image.height, y1))
        cropped = image.crop((0, y0, image.width, y1))
        buf = io.BytesIO()
        cropped.save(buf, format="PNG")
        return buf.getvalue()


def capture_sections_png(
    workbook_bytes: bytes,
    sheet_name: str,
    sections: list[Section],
    *,
    layout: SystemLayout | None = None,
    prefer_excel_com: bool = True,
) -> bytes | None:
    """
    Screenshot the System scorecard like the template EMF paste:
    month header (JAN..YE) + selected vertical category blocks in one image.
    """
    if not sections:
        return None
    if layout is None:
        layout = detect_system_layout(workbook_bytes, sheet_name=sheet_name)

    ordered = sorted(sections, key=lambda s: s.start_row)
    start_col, end_col = layout.start_col, layout.end_col
    header_start = layout.header_start_row
    header_end = layout.header_end_row
    body_start = min(s.start_row for s in ordered)
    body_end = max(s.end_row for s in ordered)

    # Final trim against the live sheet (month footer / OPPORTUNITIES).
    try:
        wb = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
        match = _find_sheet_name(list(wb.sheetnames), sheet_name) or sheet_name
        ws = wb[match]
        body_end = _trim_capture_end_row(ws, body_start, body_end, start_col, end_col)
        wb.close()
    except Exception as exc:
        logger.debug("Could not trim System capture end row: %s", exc)

    if prefer_excel_com:
        try:
            with tempfile.TemporaryDirectory(prefix="mpr_system_") as tmp:
                path = Path(tmp) / "scorecards.xlsx"
                path.write_bytes(workbook_bytes)
                png = _capture_system_sections_via_com(
                    path,
                    sheet_name,
                    header_start=header_start,
                    header_end=header_end,
                    body_start=body_start,
                    body_end=body_end,
                    start_col=start_col,
                    end_col=end_col,
                )
                with Image.open(io.BytesIO(png)) as img:
                    logger.info(
                        "System sections capture %s -> %sx%s (%s bytes) sections=%s",
                        sheet_name,
                        img.width,
                        img.height,
                        len(png),
                        [s.title for s in ordered],
                    )
                return png
        except Exception as exc:
            logger.warning(
                "System COM section capture failed (%s); falling back to range capture/render",
                exc,
            )

    # Fallback: contiguous range capture (or Pillow) without row hiding.
    return capture_range_png(
        workbook_bytes,
        sheet_name,
        header_start,
        body_end,
        start_col,
        end_col,
        prefer_excel_com=prefer_excel_com,
    )

def place_picture_on_slide(
    slide,
    png_bytes: bytes,
    *,
    left: int = DEFAULT_LEFT,
    top: int = DEFAULT_TOP,
    max_width: int = DEFAULT_WIDTH,
    max_height: int = DEFAULT_HEIGHT,
    fit: str = "fill",
    grow: float = 1.0,
    align: str = "center",
) -> None:
    """Place a PNG into the content area.

    fit=fill  -> stretch exactly into the box (matches template picture slots)
    fit=contain -> preserve aspect ratio inside the box
    grow -> expand the target box slightly (e.g. 1.08) so the scorecard reads larger
    align -> for contain: "center" (default) or "top-left" (anchor to slot origin)
    """
    from pptx.util import Emu

    if grow and grow != 1.0:
        new_w = int(max_width * grow)
        new_h = int(max_height * grow)
        left = int(left) - (new_w - int(max_width)) // 2
        top = int(top) - (new_h - int(max_height)) // 2
        max_width, max_height = new_w, new_h

    with Image.open(io.BytesIO(png_bytes)) as img:
        img_w, img_h = img.size
    if img_w <= 0 or img_h <= 0:
        return

    if fit == "contain":
        scale = min(max_width / img_w, max_height / img_h)
        width = int(img_w * scale)
        height = int(img_h * scale)
        if str(align).lower() in {"top-left", "topleft", "left", "start"}:
            left_pos = int(left)
            top_pos = int(top)
        else:
            left_pos = left + max(0, (max_width - width) // 2)
            top_pos = top + max(0, (max_height - height) // 2)
    else:
        # Fill the template picture rectangle exactly.
        width = int(max_width)
        height = int(max_height)
        left_pos = int(left)
        top_pos = int(top)

    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Emu(left_pos),
        Emu(top_pos),
        Emu(width),
        Emu(height),
    )


def _largest_picture(slide):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pictures = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
    if not pictures:
        return None
    return max(pictures, key=lambda s: int(s.width) * int(s.height))


def _remove_slide_pictures(slide) -> int:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    removed = 0
    for shape in list(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            element = shape._element
            element.getparent().remove(element)
            removed += 1
    return removed


def _remove_slide_tables_and_charts(
    slide,
    *,
    remove_tables: bool = False,
    remove_charts: bool = False,
    remove_band_labels: bool = False,
    band_top: int | None = None,
    band_bottom: int | None = None,
) -> int:
    """Remove template tables/charts/labels in a vertical band so a live screenshot can replace them."""
    if not remove_tables and not remove_charts and not remove_band_labels:
        return 0
    removed = 0
    for shape in list(slide.shapes):
        is_table = bool(getattr(shape, "has_table", False))
        is_chart = bool(getattr(shape, "has_chart", False))
        is_label = bool(
            remove_band_labels
            and getattr(shape, "has_text_frame", False)
            and not is_table
            and not is_chart
        )
        if is_table and not remove_tables:
            continue
        if is_chart and not remove_charts:
            continue
        if is_label:
            text = (shape.text_frame.text or "").strip().casefold()
            # Keep chrome titles / speaker; drop in-band labels like "System GIR".
            if not text or "speaker" in text or text in {"gir", "safety", "system gir"}:
                continue
            if int(shape.top) < 650_000:
                continue
            # Keep the template "System GIR" label above the chart.
            if "system gir" in text:
                continue
        elif not is_table and not is_chart:
            continue
        top = int(shape.top)
        bottom = top + int(shape.height)
        if band_top is not None and bottom < band_top:
            continue
        if band_bottom is not None and top > band_bottom:
            continue
        shape._element.getparent().remove(shape._element)
        removed += 1
    return removed


def resolve_sheet_name(
    workbook_bytes: bytes,
    *,
    sheet: str | None = None,
    sheet_index: int | None = None,
    sheet_match: list[str] | None = None,
    sheet_match_index: int = 0,
    available: list[str] | None = None,
) -> str:
    """Resolve a worksheet name from explicit name, index, or substring match."""
    names = available
    if names is None:
        wb = load_workbook(io.BytesIO(workbook_bytes), read_only=True, data_only=False)
        try:
            names = list(wb.sheetnames)
        finally:
            wb.close()

    if sheet:
        if sheet in names:
            return sheet
        exact = _find_sheet_name(names, sheet)
        if exact:
            return exact

    patterns = [str(p).strip().casefold() for p in (sheet_match or []) if str(p).strip()]
    if patterns:
        hits = [n for n in names if all(p in _norm_sheet_key(n) for p in patterns)]
        if not hits:
            hits = [n for n in names if any(p in _norm_sheet_key(n) for p in patterns)]
        if hits:
            idx = int(sheet_match_index or 0)
            if idx < 0:
                idx = len(hits) + idx
            if 0 <= idx < len(hits):
                return hits[idx]
            # Match found but index out of range — fall through to sheet_index.

    if sheet_index is not None:
        idx = int(sheet_index)
        if 0 <= idx < len(names):
            return names[idx]
        raise ValueError(f"sheet_index={idx} out of range for sheets {names}")

    if sheet:
        raise ValueError(f"Sheet {sheet!r} not found. Available: {names}")
    raise ValueError(f"Could not resolve sheet from index/match. Available: {names}")


def capture_sheet_png(
    workbook_bytes: bytes,
    sheet_name: str,
    *,
    prefer_excel_com: bool = True,
    max_rows: int | None = 120,
    max_cols: int | None = 40,
) -> bytes:
    """Screenshot the used range of any worksheet (entity / workings / System)."""
    wb = load_workbook(io.BytesIO(workbook_bytes), data_only=False)
    match = _find_sheet_name(list(wb.sheetnames), sheet_name)
    if match is None:
        wb.close()
        raise ValueError(f"Sheet {sheet_name!r} not found. Available: {wb.sheetnames}")
    ws = wb[match]
    min_row, max_row, min_col, max_col = _used_bounds(ws)
    wb.close()

    if max_rows is not None:
        max_row = min(max_row, min_row + int(max_rows) - 1)
    if max_cols is not None:
        max_col = min(max_col, min_col + int(max_cols) - 1)

    logger.info(
        "Capturing full sheet %s!%s",
        match,
        _range_address(min_row, max_row, min_col, max_col),
    )
    return capture_range_png(
        workbook_bytes,
        match,
        min_row,
        max_row,
        min_col,
        max_col,
        prefer_excel_com=prefer_excel_com,
    )


def apply_scorecard_screenshot(slide, data, element: dict) -> bool:
    """Replace a template picture with a live Excel screenshot (System sections or full sheet)."""
    workbook = element.get("workbook", "scorecards")
    prefer_com = bool(element.get("prefer_excel_com", True))
    replace_existing = bool(element.get("replace_existing_pictures", True))
    fit = str(element.get("fit", "fill")).lower()
    grow = float(element.get("grow", 1.0) or 1.0)
    mode = str(element.get("sections", element.get("capture", "auto"))).lower()

    try:
        workbook_bytes = data.store.workbook_bytes(workbook)
    except FileNotFoundError as exc:
        logger.warning("Scorecard screenshot skipped: %s", exc)
        return False

    available = []
    try:
        available = list(data.sheet_names(workbook))
    except Exception:
        available = []

    try:
        sheet_name = resolve_sheet_name(
            workbook_bytes,
            sheet=element.get("sheet"),
            sheet_index=element.get("sheet_index"),
            sheet_match=element.get("sheet_match"),
            sheet_match_index=int(element.get("sheet_match_index", 0) or 0),
            available=available or None,
        )
    except Exception as exc:
        logger.warning("Scorecard sheet resolve failed: %s", exc)
        return False

    explicit = element.get("range")
    png: bytes | None = None

    # Full-sheet / explicit-range capture for entity tabs, workings, etc.
    wants_full_sheet = mode in {"sheet", "full", "all", "used_range"} or (
        mode == "auto" and _norm_sheet_key(sheet_name) != "system"
    )
    if explicit:
        from openpyxl.utils.cell import range_boundaries

        min_col, min_row, max_col, max_row = range_boundaries(str(explicit))
        png = capture_range_png(
            workbook_bytes,
            sheet_name,
            min_row,
            max_row,
            min_col,
            max_col,
            prefer_excel_com=prefer_com,
        )
    elif wants_full_sheet:
        try:
            png = capture_sheet_png(
                workbook_bytes,
                sheet_name,
                prefer_excel_com=prefer_com,
                max_rows=element.get("max_rows", 120),
                max_cols=element.get("max_cols", 40),
            )
        except Exception as exc:
            logger.warning("Full-sheet capture failed for %s: %s", sheet_name, exc)
            return False
    else:
        # System tab: split into named category blocks.
        try:
            layout = _cached_system_layout(workbook_bytes, sheet_name)
        except Exception as exc:
            logger.warning("Could not detect System layout: %s", exc)
            return False

        count = int(element.get("count", 3))
        include_black = bool(element.get("include_black", False))
        match = element.get("match") or []
        fallback_indices = element.get("fallback_indices") or element.get("indices") or []
        section_mode = mode if mode in {"match", "indices", "first", "last"} else (
            "match" if match else "first"
        )
        chosen = select_sections_for_slide(
            layout.sections,
            mode=section_mode,
            count=count,
            include_black=include_black,
            match=list(match) if match else None,
            indices=[int(i) for i in fallback_indices] if fallback_indices else None,
        )
        if not chosen:
            logger.warning(
                "No System sections selected for mode=%s match=%s available=%s",
                section_mode,
                match,
                [s.title for s in layout.sections],
            )
            return False
        png = capture_sections_png(
            workbook_bytes,
            sheet_name,
            chosen,
            layout=layout,
            prefer_excel_com=prefer_com,
        )
        names = [s.title for s in chosen]
        logger.info("System screenshot sections=%s", names)
        print(
            f"\n>>> SYSTEM SCORECARD CAPTURE ({sheet_name}): "
            + " | ".join(names)
            + f"  [{len(png) if png else 0} bytes]\n"
        )

    if not png:
        print(f"\n>>> FAILED to capture screenshot for {workbook}!{sheet_name}\n")
        return False

    # Save a debug preview next to the report so capture quality is easy to review.
    try:
        with Image.open(io.BytesIO(png)) as preview:
            out_dir = Path(getattr(data.store, "base_dir", Path("."))) / "output"
            out_dir.mkdir(parents=True, exist_ok=True)
            debug_path = out_dir / (
                f"_debug_{workbook}_{_norm_sheet_key(sheet_name)}_{preview.width}x{preview.height}.png"
            )
            preview.save(debug_path)
            print(f">>> Debug preview saved: {debug_path} ({preview.width}x{preview.height})")
    except Exception as exc:
        logger.debug("Could not write debug preview: %s", exc)

    # Placement box: explicit element coords win; else reuse a template picture slot.
    has_explicit_box = any(
        element.get(key) is not None for key in ("left", "top", "max_width", "max_height")
    )
    target = _largest_picture(slide)
    if has_explicit_box:
        left = int(element.get("left", DEFAULT_LEFT))
        top = int(element.get("top", DEFAULT_TOP))
        width = int(element.get("max_width", DEFAULT_WIDTH))
        height = int(element.get("max_height", DEFAULT_HEIGHT))
    elif target is not None:
        left, top, width, height = int(target.left), int(target.top), int(target.width), int(target.height)
    else:
        left = int(element.get("left", DEFAULT_LEFT))
        top = int(element.get("top", DEFAULT_TOP))
        width = int(element.get("max_width", DEFAULT_WIDTH))
        height = int(element.get("max_height", DEFAULT_HEIGHT))
        print(
            f">>> No template picture slot found — placing at default content box "
            f"({left},{top},{width},{height})"
        )

    if replace_existing:
        removed = _remove_slide_pictures(slide)
        if removed:
            logger.info("Removed %s template picture(s) before placing live scorecard image", removed)

    # GIR (and similar) slides store old data in tables/charts — clear that band first.
    cleared = _remove_slide_tables_and_charts(
        slide,
        remove_tables=bool(element.get("remove_tables", False)),
        remove_charts=bool(element.get("remove_charts", False)),
        remove_band_labels=bool(element.get("remove_band_labels", False)),
        band_top=int(element["clear_band_top"]) if element.get("clear_band_top") is not None else top,
        band_bottom=(
            int(element["clear_band_bottom"])
            if element.get("clear_band_bottom") is not None
            else top + height
        ),
    )
    if cleared:
        logger.info("Removed %s template table/chart shape(s) before placing live image", cleared)

    place_picture_on_slide(
        slide,
        png,
        left=left,
        top=top,
        max_width=width,
        max_height=height,
        fit=fit,
        grow=grow,
    )
    logger.info(
        "Placed scorecard image from %s!%s box=(%s,%s,%s,%s) fit=%s grow=%s",
        workbook,
        sheet_name,
        left,
        top,
        width,
        height,
        fit,
        grow,
    )
    print(f">>> Placed live image on slide from {workbook}/{sheet_name}")
    return True