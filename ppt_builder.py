"""Fill the GSE MPR PowerPoint template with Excel data."""

from __future__ import annotations

import io
import logging
import math
import re
from datetime import date
from pathlib import Path

import yaml
from PIL import Image
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt
from lxml import etree

from mpr_data import MONTH_LABELS, MprData
from ppt_format import (
    clear_table_data_rows,
    clear_text_frame_content,
    set_cell_text_preserve,
    set_paragraph_text_preserve,
    set_text_frame_preserve,
)

logger = logging.getLogger(__name__)

CHART_CATEGORIES = MONTH_LABELS + ["YTD"]


def _safe_number(value: float | None) -> float | None:
    if value is None:
        return None
    try:
        number = float(value)
    except (TypeError, ValueError):
        return None
    if math.isnan(number) or math.isinf(number):
        return None
    return number


def _fmt(value: float | None, decimals: int = 2) -> str:
    number = _safe_number(value)
    if number is None:
        return ""
    if decimals == 0:
        return f"{round(number):.0f}"
    return f"{number:.{decimals}f}"


def _fmt_diff(actual: float | None, goal: float | None, decimals: int = 2) -> str:
    actual_n = _safe_number(actual)
    goal_n = _safe_number(goal)
    if actual_n is None or goal_n is None:
        return ""
    diff = actual_n - goal_n
    if diff < 0:
        return f"({abs(diff):.{decimals}f})"
    return _fmt(diff, decimals)


def _iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)


def _remove_shape(shape) -> None:
    element = shape._element
    element.getparent().remove(element)


def _bring_to_front(slide, shape) -> None:
    """Move a shape to the top of the z-order so it is clickable."""
    tree = slide.shapes._spTree
    element = shape._element
    tree.remove(element)
    tree.append(element)


def _is_title_shape(slide, shape) -> bool:
    if slide.shapes.title is not None and shape.shape_id == slide.shapes.title.shape_id:
        return True
    if not shape.has_text_frame:
        return False
    text = shape.text_frame.text.strip().lower()
    if not text:
        return False
    if len(text) < 80 and shape.top < Pt(120):
        return True
    return False


def _replace_text(text: str, data: MprData) -> str:
    month_label = data.report_month_label()
    month_short = data.report_month_short()
    updated = text
    updated = re.sub(
        r"\b(January|February|March|April|May|June|July|August|September|October|November|December)\s+\d{4}\b",
        month_label,
        updated,
        flags=re.IGNORECASE,
    )
    updated = re.sub(r"\b[A-Za-z]{3}'\d{2}\b", month_short, updated)
    return updated


def replace_month_tokens_on_slide(slide, data: MprData) -> None:
    """Update month/year tokens on a single slide only."""
    for shape in _iter_all_shapes(slide.shapes):
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            original = paragraph.text
            if not original.strip():
                continue
            updated = _replace_text(original, data)
            if updated != original:
                set_paragraph_text_preserve(paragraph, updated)


# Planned Discussion (slide 2) — static agenda matching the GSE MPR template.
AGENDA_TOPICS: list[tuple[str, str]] = [
    ("Safety", "10 min"),
    ("People", "10 min"),
    ("Finance", "10 min"),
    ("Customer Experience", "10 min"),
    ("North (M)", "10 min"),
    ("South (M)", "10 min"),
    ("Stationary", "15 min"),
    ("Galley", "10 min"),
    ("Closing", "3 min"),
]

def _agenda_chevrons(slide):
    """Tall off-page-connector shapes that form the nine topic columns."""
    chevrons = [
        shape
        for shape in slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        and "off-page" in (shape.name or "").casefold()
        and int(shape.height) > 2_000_000
    ]
    chevrons.sort(key=lambda s: int(s.left))
    return chevrons


def _agenda_divider_lines(slide, chevron) -> list[int]:
    """Y positions of horizontal white divider lines inside one chevron."""
    left = int(chevron.left)
    right = left + int(chevron.width)
    tops: list[int] = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.LINE:
            continue
        mid = int(shape.left) + int(shape.width) // 2
        top = int(shape.top)
        if left - 80_000 <= mid <= right + 80_000 and 2_100_000 < top < 4_100_000:
            tops.append(top)
    return sorted(set(tops))


def _style_agenda_note_box(box) -> None:
    """Empty editable note slot between white divider lines (matches template)."""
    tf = box.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = None
    except Exception:
        pass
    try:
        box.text_frame._txBody.bodyPr.set("anchor", "ctr")  # type: ignore[attr-defined]
    except Exception:
        pass

    # No box outline — the white divider lines are the only borders.
    try:
        sp_pr = box._element.spPr
        for old in list(sp_pr.findall(qn("a:ln"))):
            sp_pr.remove(old)
        ln = etree.SubElement(sp_pr, qn("a:ln"))
        ln.set("w", "0")
        etree.SubElement(ln, qn("a:noFill"))
    except Exception:
        pass

    if not tf.paragraphs:
        return
    # Clear leftover template/placeholder text so slots look empty.
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    if para.runs:
        para.runs[0].text = ""
        for run in para.runs[1:]:
            run.text = ""
        run = para.runs[0]
    else:
        run = para.add_run()
        run.text = ""
    run.font.size = Pt(11)
    run.font.bold = False
    run.font.color.rgb = RGBColor(255, 255, 255)
    for extra in tf.paragraphs[1:]:
        for run in extra.runs:
            run.text = ""


def _ensure_agenda_line_textboxes(slide) -> int:
    """Put one empty clickable note text box in each white-line gap of every chevron."""
    built = 0
    chevrons = _agenda_chevrons(slide)

    # Remove oversized body boxes that span multiple slots (incl. leftover text).
    for chevron in chevrons:
        lines = _agenda_divider_lines(slide, chevron)
        if len(lines) < 2:
            continue
        avg_slot = (lines[-1] - lines[0]) / max(1, len(lines) - 1)
        col_left = int(chevron.left)
        col_right = col_left + int(chevron.width)
        for shape in list(slide.shapes):
            if not shape.has_text_frame:
                continue
            mid_x = int(shape.left) + int(shape.width) // 2
            if not (col_left - 40_000 <= mid_x <= col_right + 40_000):
                continue
            top = int(shape.top)
            if top < 2_150_000 or top > 3_900_000:
                continue
            if int(shape.height) > avg_slot * 1.4:
                _remove_shape(shape)

    for chevron in chevrons:
        lines = _agenda_divider_lines(slide, chevron)
        if len(lines) < 2:
            body_top = int(chevron.top) + int(chevron.height) * 28 // 100
            body_bottom = int(chevron.top) + int(chevron.height) * 88 // 100
            step = max(1, (body_bottom - body_top) // 3)
            lines = [body_top + i * step for i in range(4)]
            lines[-1] = body_bottom

        col_left = int(chevron.left) + 50_000
        col_width = max(200_000, int(chevron.width) - 100_000)

        for idx in range(len(lines) - 1):
            slot_top = lines[idx] + 25_000
            slot_bottom = lines[idx + 1] - 25_000
            if slot_bottom - slot_top < 90_000:
                continue

            existing = None
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                mid_x = int(shape.left) + int(shape.width) // 2
                mid_y = int(shape.top) + int(shape.height) // 2
                if not (col_left - 60_000 <= mid_x <= col_left + col_width + 60_000):
                    continue
                if slot_top <= mid_y <= slot_bottom and int(shape.height) < (slot_bottom - slot_top) * 1.6:
                    existing = shape
                    break

            if existing is None:
                existing = slide.shapes.add_textbox(
                    Emu(col_left),
                    Emu(slot_top),
                    Emu(col_width),
                    Emu(slot_bottom - slot_top),
                )
                existing.name = f"AgendaNote_{built + 1}"
                built += 1
            else:
                # Resize into the slot so it sits cleanly between the white lines.
                existing.left = Emu(col_left)
                existing.top = Emu(slot_top)
                existing.width = Emu(col_width)
                existing.height = Emu(slot_bottom - slot_top)
                if not (existing.name or "").startswith("AgendaNote_"):
                    existing.name = f"AgendaNote_{built + 1}"
                    built += 1

            _style_agenda_note_box(existing)
            _bring_to_front(slide, existing)

    return built


def fill_agenda_slide(slide) -> None:
    """Restore Planned Discussion titles/times; clear note slots between white lines."""
    title_boxes = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame
        and 1_700_000 <= int(shape.top) <= 2_100_000
        and int(shape.width) > 500_000
    ]
    time_ovals = [
        shape
        for shape in slide.shapes
        if shape.has_text_frame
        and "oval" in (shape.name or "").lower()
        and int(shape.top) > 4_700_000
    ]
    title_boxes.sort(key=lambda s: int(s.left))
    time_ovals.sort(key=lambda s: int(s.left))

    for idx, (topic, minutes) in enumerate(AGENDA_TOPICS):
        if idx < len(title_boxes):
            set_text_frame_preserve(title_boxes[idx].text_frame, topic)
            _bring_to_front(slide, title_boxes[idx])
        if idx < len(time_ovals):
            set_text_frame_preserve(time_ovals[idx].text_frame, minutes)

    built = _ensure_agenda_line_textboxes(slide)
    logger.info(
        "Planned Discussion ready: %s topics, %s empty editable note slots between white lines",
        len(AGENDA_TOPICS),
        built,
    )


def _build_chart_data(
    actuals: list[float | None],
    goal: float | None,
    existing_chart,
    *,
    ytd: float | None = None,
    prior_year: list[float | None] | None = None,
    prior_ytd: float | None = None,
) -> CategoryChartData:
    """Build chart categories Jan-Dec + YTD.

    GIR YTD must come from Workings (weighted rate) — never a simple monthly average.
    """
    if ytd is None:
        # Prefer an explicit 13th value if caller already included it.
        if len(actuals) >= 13 and _safe_number(actuals[12]) is not None:
            ytd = actuals[12]
    full_actuals = list(actuals[:12]) + [ytd]

    prior_full = None
    if prior_year is not None:
        prior_full = list(prior_year[:12])
        while len(prior_full) < 12:
            prior_full.append(None)
        prior_full.append(prior_ytd)

    chart_data = CategoryChartData()
    chart_data.categories = CHART_CATEGORIES

    if existing_chart and existing_chart.series:
        for series in existing_chart.series:
            name = series.name or ""
            lower = name.lower()
            if lower in ("actual", "current year") or "actual" in lower or "current" in lower:
                chart_data.add_series(name or "Current Year", tuple(full_actuals))
            elif "goal" in lower:
                chart_data.add_series(
                    name or "Goal",
                    tuple([goal] * 13 if _safe_number(goal) is not None else [None] * 13),
                )
            elif prior_full is not None and ("p1y" in lower or "prior" in lower):
                chart_data.add_series(name or "P1Y", tuple(prior_full))
            else:
                existing = list(series.values) if series.values is not None else []
                while len(existing) < 13:
                    existing.append(None)
                chart_data.add_series(name, tuple(existing[:13]))
    else:
        chart_data.add_series("Current Year", tuple(full_actuals))
        if _safe_number(goal) is not None:
            chart_data.add_series("Goal", tuple([goal] * 13))
        if prior_full is not None:
            chart_data.add_series("P1Y", tuple(prior_full))

    return chart_data


def _update_chart_series(
    slide,
    shape,
    actuals: list[float | None],
    goal: float | None,
    *,
    ytd: float | None = None,
    prior_year: list[float | None] | None = None,
    prior_ytd: float | None = None,
) -> bool:
    if not shape.has_chart:
        return False
    if not any(_safe_number(v) is not None for v in actuals):
        return False

    chart_data = _build_chart_data(
        actuals,
        goal,
        shape.chart,
        ytd=ytd,
        prior_year=prior_year,
        prior_ytd=prior_ytd,
    )
    chart_title = ""
    chart_type = XL_CHART_TYPE.LINE_MARKERS
    if shape.has_chart:
        try:
            chart_type = shape.chart.chart_type
            if shape.chart.has_title:
                chart_title = shape.chart.chart_title.text_frame.text
        except Exception:
            pass
    try:
        shape.chart.replace_data(chart_data)
        return True
    except (ValueError, AttributeError) as exc:
        logger.info("Recreating linked/external chart: %s", exc)
        left, top, width, height = shape.left, shape.top, shape.width, shape.height
        _remove_shape(shape)
        new_shape = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
        if chart_title and new_shape.chart.has_title:
            set_text_frame_preserve(new_shape.chart.chart_title.text_frame, chart_title)
        return True


def update_gir_chart_from_series(
    slide,
    *,
    monthly: list[float | None],
    goal: float | None,
    ytd: float | None = None,
    prior_year: list[float | None] | None = None,
    prior_ytd: float | None = None,
) -> bool:
    """Update the System GIR chart from Workings series values."""
    updated = False
    for shape in list(slide.shapes):
        if shape.has_chart:
            if _update_chart_series(
                slide,
                shape,
                monthly,
                goal,
                ytd=ytd,
                prior_year=prior_year,
                prior_ytd=prior_ytd,
            ):
                updated = True
            else:
                _remove_shape(shape)
    return updated


def _clear_non_title_content(slide, *, keep_title: bool = True) -> None:
    for shape in list(slide.shapes):
        if keep_title and _is_title_shape(slide, shape):
            continue
        if shape.has_table:
            clear_table_data_rows(shape.table, header_rows=1)
        elif shape.has_chart:
            _remove_shape(shape)
        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            _remove_shape(shape)
        elif shape.has_text_frame:
            clear_text_frame_content(shape.text_frame)


def _fill_table_from_dataframe(table, df, *, header_rows: int = 1) -> bool:
    if df.empty:
        return False
    max_rows = len(table.rows) - header_rows
    max_cols = len(table.columns)
    wrote = False
    for r in range(min(max_rows, len(df.index))):
        for c in range(min(max_cols, len(df.columns))):
            val = df.iat[r, c]
            if val is None or (isinstance(val, float) and math.isnan(val)):
                text = ""
            else:
                text = str(val).strip()
            if text:
                set_cell_text_preserve(table.cell(header_rows + r, c), text)
                wrote = True
    return wrote


def update_gir_tables(slide, data: MprData, config: dict, workbook: str) -> None:
    """Fill native GIR template tables from Actuals (and Injury Breakdown from Workings when clean)."""
    kpi = config.get("kpi_mappings", {})
    gir_name = kpi.get("gir", "GIR")
    injury_name = kpi.get("injury_count", "Injury Count")

    mtd = data.kpi_value([gir_name], month=data.month, workbook=workbook)
    ytd_actual = data.ytd_value([gir_name], workbook=workbook)
    monthly = data.monthly_series([gir_name], through_month=data.month, workbook=workbook)
    injury_mtd = data.kpi_value([injury_name], month=data.month, workbook=workbook)
    injury_monthly = data.monthly_series([injury_name], through_month=data.month, workbook=workbook)
    injury_ytd = data.ytd_value([injury_name], workbook=workbook)
    yo1 = data.kpi_value([gir_name], month=data.month, workbook=workbook, year=data.year - 1)
    yo2 = data.kpi_value([gir_name], month=data.month, workbook=workbook, year=data.year - 2)

    injury_rows = None
    try:
        from gir_panels import read_injury_breakdown_rows

        injury_rows = read_injury_breakdown_rows(data)
    except Exception as exc:
        logger.info("GIR Injury Breakdown from workings unavailable: %s", exc)

    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        label = table.cell(0, 0).text.strip().lower()

        if label == "metric" and table.cell(2, 0).text.strip().upper() == "GIR":
            set_cell_text_preserve(table.cell(2, 1), _fmt(mtd.actual))
            set_cell_text_preserve(table.cell(2, 2), _fmt(mtd.goal))
            set_cell_text_preserve(table.cell(2, 3), _fmt_diff(mtd.actual, mtd.goal))
            set_cell_text_preserve(table.cell(2, 4), _fmt(ytd_actual))
            set_cell_text_preserve(table.cell(2, 5), _fmt(mtd.goal))
            set_cell_text_preserve(table.cell(2, 6), _fmt_diff(ytd_actual, mtd.goal))

        if "injury breakdown" in label and injury_rows:
            _fill_gir_injury_breakdown_table(table, injury_rows)

        if "actual:" in table.cell(1, 0).text.lower() or data.report_month_short().lower() in table.cell(0, 0).text.lower():
            set_cell_text_preserve(table.cell(0, 0), data.report_month_short())
            set_cell_text_preserve(table.cell(1, 1), _fmt(mtd.actual))
            if len(table.rows) > 2:
                set_cell_text_preserve(table.cell(2, 1), _fmt(mtd.goal))
            for row_idx in range(len(table.rows)):
                left = table.cell(row_idx, 0).text.strip().casefold()
                if left.startswith("yo1y") and _safe_number(yo1.actual) is not None:
                    set_cell_text_preserve(table.cell(row_idx, 1), _fmt(yo1.actual))
                elif left.startswith("yo2y") and _safe_number(yo2.actual) is not None:
                    set_cell_text_preserve(table.cell(row_idx, 1), _fmt(yo2.actual))

        if label == "recordable":
            for col_idx, month_num in enumerate(range(1, 13), start=1):
                if month_num <= data.month:
                    set_cell_text_preserve(table.cell(1, col_idx), _fmt(monthly[month_num - 1]))
                else:
                    set_cell_text_preserve(table.cell(1, col_idx), "")
            set_cell_text_preserve(table.cell(1, 13), _fmt(ytd_actual))
            set_cell_text_preserve(table.cell(1, 14), _fmt_diff(ytd_actual, mtd.goal))
            if len(table.rows) > 2 and "injury" in table.cell(2, 0).text.lower():
                for col_idx, month_num in enumerate(range(1, 13), start=1):
                    if month_num <= data.month and _safe_number(injury_monthly[month_num - 1]) is not None:
                        set_cell_text_preserve(table.cell(2, col_idx), _fmt(injury_monthly[month_num - 1], decimals=0))
                    elif month_num > data.month:
                        set_cell_text_preserve(table.cell(2, col_idx), "")
                if _safe_number(injury_ytd) is not None:
                    set_cell_text_preserve(table.cell(2, 13), _fmt(injury_ytd, decimals=0))
                else:
                    set_cell_text_preserve(table.cell(2, 13), "")


def _fill_gir_injury_breakdown_table(table, rows: list[list[str]]) -> bool:
    """Write Injury Breakdown values into the existing template table without resizing it."""
    if not rows:
        return False
    wrote = False
    # Keep row 0 title / row 1 headers; write data from row 2 onward when shapes match.
    start_row = 2 if len(table.rows) > 2 else 1
    max_cols = len(table.columns)
    for r_idx, row in enumerate(rows):
        ppt_row = start_row + r_idx
        if ppt_row >= len(table.rows):
            break
        for c_idx in range(min(max_cols, len(row))):
            value = row[c_idx]
            if value is None:
                continue
            set_cell_text_preserve(table.cell(ppt_row, c_idx), str(value))
            wrote = True
    if wrote:
        logger.info("Filled GIR Injury Breakdown table (%s data row(s))", min(len(rows), len(table.rows) - start_row))
    return wrote

def update_gir_charts(slide, data: MprData, config: dict, workbook: str) -> None:
    gir_name = config.get("kpi_mappings", {}).get("gir", "GIR")
    monthly = data.monthly_series([gir_name], through_month=data.month, workbook=workbook)
    goal = data.kpi_value([gir_name], month=data.month, workbook=workbook).goal
    ytd = data.ytd_value([gir_name], workbook=workbook)
    update_gir_chart_from_series(slide, monthly=monthly, goal=goal, ytd=ytd)

def update_people_table(slide, data: MprData, rows_cfg: list[dict], workbook: str) -> None:
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        for row_idx in range(2, len(table.rows)):
            label = table.cell(row_idx, 0).text.strip().upper()
            for row_def in rows_cfg:
                if row_def["label"] in label:
                    patterns = row_def["patterns"]
                    mtd = data.kpi_value(patterns, month=data.month, workbook=workbook)
                    ytd = data.ytd_value(patterns, workbook=workbook)
                    set_cell_text_preserve(
                        table.cell(row_idx, 1),
                        _fmt(mtd.actual, 0) if _safe_number(mtd.actual) is not None else "",
                    )
                    set_cell_text_preserve(
                        table.cell(row_idx, 3),
                        _fmt(ytd, 0) if _safe_number(ytd) is not None else "",
                    )
                    if _safe_number(mtd.goal) is not None and _safe_number(ytd) is not None:
                        set_cell_text_preserve(table.cell(row_idx, 4), _fmt_diff(ytd, mtd.goal, 0))
                    else:
                        set_cell_text_preserve(table.cell(row_idx, 4), "")
                    break


def update_people_charts(slide, data: MprData, charts_cfg: list[dict], workbook: str) -> None:
    for shape in list(slide.shapes):
        if not shape.has_chart:
            continue
        title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else ""
        for chart_def in charts_cfg:
            if chart_def["title"].lower() in title.lower():
                monthly = data.monthly_series(chart_def["patterns"], through_month=data.month, workbook=workbook)
                goal = data.kpi_value(chart_def["patterns"], month=data.month, workbook=workbook).goal
                if any(_safe_number(v) is not None for v in monthly):
                    _update_chart_series(slide, shape, monthly, goal)
                else:
                    _remove_shape(shape)
                break


def update_chart_slide_mapped(slide, data: MprData, charts_cfg: list[dict], workbook: str) -> None:
    charts = [s for s in list(slide.shapes) if s.has_chart]
    for shape in charts:
        title = shape.chart.chart_title.text_frame.text if shape.chart.has_title else shape.name
        matched = False
        for chart_def in charts_cfg:
            chart_title = chart_def.get("title", "")
            if chart_title and chart_title.lower() not in title.lower():
                continue
            if not chart_title and len(charts_cfg) == 1:
                pass
            elif not chart_title:
                continue
            monthly = data.monthly_series(chart_def["patterns"], through_month=data.month, workbook=workbook)
            goal = data.kpi_value(chart_def["patterns"], month=data.month, workbook=workbook).goal
            if any(_safe_number(v) is not None for v in monthly):
                _update_chart_series(slide, shape, monthly, goal)
                logger.info("Updated chart %r", title)
            else:
                _remove_shape(shape)
            matched = True
            break
        if not matched and charts_cfg:
            _remove_shape(shape)


def apply_scorecard_sheet(slide, data: MprData, element: dict) -> bool:
    workbook = element["workbook"]
    sheet_index = element.get("sheet_index", 0)
    sheet_name = element.get("sheet")
    if sheet_name:
        df = data.read_sheet(workbook, sheet_name)
    else:
        names = data.sheet_names(workbook)
        if sheet_index >= len(names):
            return False
        df = data.read_sheet(workbook, names[sheet_index])

    if df.empty:
        return False

    for shape in slide.shapes:
        if shape.has_table:
            if _fill_table_from_dataframe(shape.table, df.head(len(shape.table.rows) - 1)):
                return True
    return False


def apply_scorecard_screenshot_element(slide, data: MprData, element: dict) -> bool:
    from scorecard_screenshots import apply_scorecard_screenshot

    return apply_scorecard_screenshot(slide, data, element)


def apply_workings_table(slide, data: MprData, element: dict) -> bool:
    workbook = element["workbook"]
    sheet = element.get("sheet")
    if sheet:
        df = data.read_sheet(workbook, sheet)
    else:
        names = data.sheet_names(workbook)
        if not names:
            return False
        df = data.read_sheet(workbook, names[0])
    if df.empty:
        return False
    for shape in slide.shapes:
        if shape.has_table:
            if _fill_table_from_dataframe(shape.table, df):
                return True
    return False


def _apply_element(slide, data: MprData, config: dict, element: dict) -> None:
    etype = element["type"]
    workbook = element.get("workbook", "actuals")
    optional = bool(element.get("optional", False))

    if etype == "month_tokens":
        replace_month_tokens_on_slide(slide, data)
    elif etype == "agenda":
        fill_agenda_slide(slide)
    elif etype == "gir_tables":
        update_gir_tables(slide, data, config, workbook)
    elif etype == "gir_charts":
        update_gir_charts(slide, data, config, workbook)
    elif etype == "gir_workings_panels":
        from gir_panels import apply_gir_workings_panels

        if not apply_gir_workings_panels(slide, data, element) and not optional:
            logger.warning("No GIR workings panels for element %s", element)
    elif etype == "gir_workings_native":
        from gir_workings import apply_gir_workings_native

        if not apply_gir_workings_native(slide, data, element) and not optional:
            logger.warning("No Workings GIR data for element %s", element)
    elif etype == "ea_asap_workings_panels":
        from ea_asap_panels import apply_ea_asap_workings_panels

        if not apply_ea_asap_workings_panels(slide, data, element) and not optional:
            logger.warning("No Workings EAC/ASAP panels for element %s", element)
    elif etype == "people_workings_panels":
        from people_workings import apply_people_workings_panels

        if not apply_people_workings_panels(slide, data, element) and not optional:
            logger.warning("No Workings PEOPLE panels for element %s", element)
    elif etype == "finance_workings_panels":
        from finance_workings import apply_finance_workings_panels

        if not apply_finance_workings_panels(slide, data, element) and not optional:
            logger.warning("No Workings FINANCE panels for element %s", element)
    elif etype == "pmi_workings_panels":
        from pmi_workings import apply_pmi_workings_panels

        if not apply_pmi_workings_panels(slide, data, element) and not optional:
            logger.warning("No Workings PMI panels for element %s", element)
    elif etype == "isr_workings_panels":
        from isr_workings import apply_isr_workings_panels

        if not apply_isr_workings_panels(slide, data, element) and not optional:
            logger.warning("No Workings ISR panels for element %s", element)
    elif etype == "clear_narrative":
        from gir_panels import clear_leading_action_narrative

        n = clear_leading_action_narrative(slide)
        print(f">>> Cleared Leading Issues / Action Plan text ({n} box(es)) on slide")
    elif etype == "people_table":
        update_people_table(slide, data, element.get("rows", []), workbook)
    elif etype == "people_charts":
        update_people_charts(slide, data, element.get("charts", []), workbook)
    elif etype == "chart_slide":
        update_chart_slide_mapped(slide, data, element.get("charts", []), workbook)
    elif etype == "scorecard_sheet":
        if not apply_scorecard_sheet(slide, data, element) and not optional:
            logger.warning("No scorecard data for slide element %s", element)
    elif etype == "scorecard_screenshot":
        if not apply_scorecard_screenshot_element(slide, data, element) and not optional:
            logger.warning("No scorecard screenshot for element %s", element)
    elif etype == "workings_table":
        if not apply_workings_table(slide, data, element) and not optional:
            logger.warning("No workings data for slide element %s", element)


def _apply_slide_spec(slide, data: MprData, config: dict, slide_spec: dict) -> None:
    """Apply mapped updates while preserving the rest of the template slide."""
    for element in slide_spec.get("elements", []):
        _apply_element(slide, data, config, element)


def load_template_map(base_dir: Path) -> list[dict]:
    path = base_dir / "template_map.yaml"
    if not path.exists():
        return []
    with open(path, "r", encoding="utf-8") as handle:
        payload = yaml.safe_load(handle) or {}
    return payload.get("slides", [])


def _load_template(config: dict, base_dir: Path) -> Presentation:
    template_rel = config["powerpoint"]["template_path"]
    cached = config.get("_sharepoint_files", {}).get(template_rel)
    if cached:
        logger.info("Using PowerPoint template from SharePoint memory (%s bytes)", len(cached))
        return Presentation(io.BytesIO(cached))

    from sharepoint_live import get_cached_file

    cached = get_cached_file(template_rel)
    if cached:
        return Presentation(io.BytesIO(cached))

    template_path = base_dir / template_rel
    if not template_path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")
    return Presentation(str(template_path))


def _verify_output(prs: Presentation) -> None:
    """Log clear pass/fail checks so a bad build is obvious in the console."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    if len(prs.slides) < 4:
        logger.error("Output has only %s slides — template looks incomplete", len(prs.slides))
        return

    agenda = prs.slides[1]
    note_slots = [
        s
        for s in agenda.shapes
        if s.has_text_frame
        and 2_150_000 <= int(s.top) <= 3_900_000
        and int(s.height) < 500_000
        and (
            (s.name or "").startswith("AgendaNote_")
            or not (s.text_frame.text or "").strip()
        )
    ]
    leftover_text = sum(
        1
        for s in agenda.shapes
        if s.has_text_frame
        and 2_150_000 <= int(s.top) <= 3_900_000
        and int(s.height) < 500_000
        and (s.text_frame.text or "").strip()
    )
    print(f"\nVERIFY slide 2 (agenda): {len(note_slots)} empty editable note slots")
    if len(note_slots) < 20:
        logger.warning("Agenda note slots look incomplete (%s) — expected ~27", len(note_slots))
    if leftover_text:
        logger.warning("Agenda still has %s note-band shape(s) with leftover text", leftover_text)

    for idx, label in ((2, "slide 3 System"), (3, "slide 4 System")):
        slide = prs.slides[idx]
        pics = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE]
        print(f"VERIFY {label}: {len(pics)} picture(s)")
        if not pics:
            logger.error(
                "%s has NO screenshot picture — capture failed.",
                label,
            )
            continue
        # Blob size + pixel size catches blank/tiny COM captures.
        for pic in pics:
            try:
                blob = pic.image.blob
                with Image.open(io.BytesIO(blob)) as img:
                    sample = img.resize((max(1, img.width // 10), max(1, img.height // 10)))
                    pixels = list(sample.convert("RGB").getdata())
                    whiteish = (
                        sum(1 for r, g, b in pixels if r >= 245 and g >= 245 and b >= 245) / len(pixels)
                        if pixels
                        else 1.0
                    )
                    print(
                        f"VERIFY {label} image: {len(blob)} bytes, {img.width}x{img.height}, "
                        f"white={whiteish:.0%}"
                    )
                    if (
                        len(blob) < 20_000
                        or img.width < 400
                        or img.height < 200
                        or whiteish >= 0.90
                    ):
                        logger.error(
                            "%s screenshot looks empty/blank (%s bytes, %sx%s, white=%.0f%%)",
                            label,
                            len(blob),
                            img.width,
                            img.height,
                            whiteish * 100,
                        )
            except Exception as exc:
                logger.warning("%s could not inspect image blob: %s", label, exc)

    # GIR should keep native chart + tables filled from Workings!GIR.
    gir = prs.slides[4]
    gir_charts = sum(1 for s in gir.shapes if getattr(s, "has_chart", False))
    gir_tables = sum(1 for s in gir.shapes if getattr(s, "has_table", False))
    gir_pics = sum(1 for s in gir.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
    print(f"VERIFY slide 5 GIR native: {gir_tables} table(s), {gir_charts} chart(s), {gir_pics} picture(s)")
    if gir_tables < 3:
        logger.warning("Slide 5 GIR expected native tables, found %s", gir_tables)
    if gir_charts < 1:
        logger.warning("Slide 5 GIR expected a native chart, found %s", gir_charts)
    if gir_pics > 1:
        logger.warning("Slide 5 GIR has unexpected pictures (%s) — prefer native fills", gir_pics)
    for shape in gir.shapes:
        if not getattr(shape, "has_table", False):
            continue
        table = shape.table
        label = (table.cell(0, 0).text or "").strip().casefold()
        if label == "recordable" and len(table.rows) > 2:
            inj = (table.cell(2, 2).text or "").strip()  # Feb injury count
            if inj in ("", "0"):
                logger.warning(
                    "Slide 5 GIR Injury Count still looks empty/zero (expected Workings values)"
                )
            break
    narrative_left = [
        (s.text_frame.text or "").strip()
        for s in gir.shapes
        if getattr(s, "has_text_frame", False) and int(s.top) >= 4_900_000 and int(s.height) > 400_000
    ]
    leftover = [t for t in narrative_left if t]
    if leftover:
        logger.warning("Slide 5 GIR narrative boxes still have text: %s", leftover[:2])
    else:
        print("VERIFY slide 5 GIR narrative: Leading Issues / Action Plan boxes are empty")
    # Slide 6 EA/ASAP: exactly one table screenshot in the template OLE slot.
    ea = prs.slides[5]
    ea_pics = [
        s
        for s in ea.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and int(s.top) < 6_000_000
    ]
    ea_ole = [
        s
        for s in ea.shapes
        if s.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
    ]
    print(f"VERIFY slide 6 EA/ASAP content pictures: {len(ea_pics)} (OLE left: {len(ea_ole)})")
    if len(ea_pics) != 1:
        logger.warning(
            "Slide 6 EA/ASAP expected exactly 1 content screenshot, found %s",
            len(ea_pics),
        )
    if ea_ole:
        logger.warning("Slide 6 EA/ASAP still has %s embedded OLE object(s)", len(ea_ole))
    if ea_pics:
        pic = ea_pics[0]
        # Should fill the enlarged left content slot (bigger than original OLE).
        if int(pic.width) < 7_500_000 or int(pic.height) < 4_000_000:
            logger.warning(
                "Slide 6 EA/ASAP image smaller than expected slot (W=%s H=%s)",
                int(pic.width),
                int(pic.height),
            )
        right = int(pic.left) + int(pic.width)
        bottom = int(pic.top) + int(pic.height)
        if right > 8_400_000:
            logger.warning("Slide 6 EA/ASAP image may overlap Leading Issues (right=%s)", right)
        if bottom > 6_200_000:
            logger.warning("Slide 6 EA/ASAP image may overlap footer logo (bottom=%s)", bottom)
    ea_textboxes = [
        s
        for s in ea.shapes
        if getattr(s, "has_text_frame", False)
        and "textbox" in (s.name or "").casefold()
        and int(s.height) > 500_000
        and int(s.top) > 1_000_000
    ]
    ea_leftover = [
        (s.text_frame.text or "").strip()
        for s in ea_textboxes
        if (s.text_frame.text or "").strip()
        and "leading" not in (s.text_frame.text or "").casefold()
        and not (s.text_frame.text or "").casefold().startswith("action plan")
    ]
    print(f"VERIFY slide 6 EA narrative body text boxes: {len(ea_textboxes)}")
    print(f"VERIFY slide 6 EA narrative leftover bodies: {len(ea_leftover)}")
    if len(ea_textboxes) < 2:
        logger.warning("Slide 6 EA expected empty Leading Issues / Action Plan text boxes, found %s", len(ea_textboxes))
    if ea_leftover:
        logger.warning("Slide 6 EA narrative still has text: %s", ea_leftover[:2])
    else:
        print("VERIFY slide 6 EA narrative: Leading Issues / Action Plan boxes are empty")

    # Slide 7 People: table screenshot + 3 Excel graph screenshots on the right.
    people = prs.slides[6]
    people_tables = sum(1 for s in people.shapes if getattr(s, "has_table", False))
    people_charts = sum(1 for s in people.shapes if getattr(s, "has_chart", False))
    people_pics = [
        s
        for s in people.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and int(s.top) < 6_000_000
    ]
    right_pics = [
        s
        for s in people.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE
        and int(s.left) >= 7_500_000
        and int(s.top) < 6_200_000
    ]
    print(
        f"VERIFY slide 7 People: {len(people_pics)} content picture(s), "
        f"{len(right_pics)} right-side graph screenshot(s), "
        f"{people_tables} table(s), {people_charts} native chart(s)"
    )
    if len(people_pics) < 1:
        logger.warning("Slide 7 People expected Workings PEOPLE table screenshot, found none")
    if people_tables:
        logger.warning("Slide 7 People still has %s native table(s)", people_tables)
    if people_charts:
        logger.warning(
            "Slide 7 People still has %s native chart(s) — expected Excel graph screenshots only",
            people_charts,
        )
    if len(right_pics) < 3:
        logger.warning(
            "Slide 7 People expected 3 Excel graph screenshots on the right "
            "(Leadership Engagement, Psychological Safety, Accountability), found %s",
            len(right_pics),
        )
    else:
        print(
            "VERIFY slide 7 People: 3 Excel graph screenshots on the right "
            "(Leadership Engagement / Psychological Safety / Accountability)"
        )

    # Slide 8 Finance: Regions/BUDGET/OVERTIME/TOTAL HOURS screenshot from Workings!FINANCE.
    finance = prs.slides[7]
    finance_pics = [
        s
        for s in finance.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and int(s.top) < 5_500_000
    ]
    print(f"VERIFY slide 8 Finance: {len(finance_pics)} content picture(s)")
    if len(finance_pics) < 1:
        logger.warning(
            "Slide 8 Finance expected Workings FINANCE Regions/BUDGET/OT/Hours screenshot, found none"
        )
    else:
        print(
            "VERIFY slide 8 Finance: Regions/BUDGET/OVERTIME/TOTAL HOURS screenshot present"
        )

    # Slides 9-10 Finance comments: Leading Issues / Action Plan body boxes empty but present.
    for slide_idx, label in ((8, "slide 9"), (9, "slide 10")):
        slide = prs.slides[slide_idx]
        bodies = [
            s
            for s in slide.shapes
            if getattr(s, "has_text_frame", False)
            and "textbox" in (s.name or "").casefold()
            and int(s.height) > 500_000
            and int(s.top) > 1_000_000
        ]
        leftover = [
            (s.text_frame.text or "").strip()
            for s in bodies
            if (s.text_frame.text or "").strip()
            and "leading" not in (s.text_frame.text or "").casefold()
            and not (s.text_frame.text or "").casefold().startswith("action plan")
            and "budget" not in (s.text_frame.text or "").casefold()
        ]
        print(f"VERIFY {label} Finance comments body text boxes: {len(bodies)}")
        if len(bodies) < 2:
            logger.warning(
                "%s Finance comments expected empty editable narrative text boxes, found %s",
                label,
                len(bodies),
            )
        if leftover:
            logger.warning("%s Finance comments narrative still has text: %s", label, leftover[:2])
        else:
            print(f"VERIFY {label} Finance comments: Leading Issues / Action Plan boxes are empty")

    # Slide 11 PMI: table screenshot + 2 Excel graph screenshots; narrative empty.
    pmi = prs.slides[10]
    pmi_ole = sum(
        1
        for s in pmi.shapes
        if s.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
    )
    pmi_charts = sum(1 for s in pmi.shapes if getattr(s, "has_chart", False))
    pmi_pics = [
        s
        for s in pmi.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and int(s.top) < 6_800_000
    ]
    pmi_table_pics = [s for s in pmi_pics if int(s.top) < 5_000_000]
    pmi_chart_pics = [s for s in pmi_pics if int(s.top) >= 5_000_000]
    print(
        f"VERIFY slide 11 PMI: {len(pmi_table_pics)} table pic(s), "
        f"{len(pmi_chart_pics)} graph screenshot(s), "
        f"{pmi_ole} OLE, {pmi_charts} native chart(s)"
    )
    if len(pmi_table_pics) < 1:
        logger.warning(
            "Slide 11 PMI expected Workings Regions MOTORIZED/STATIONARY + NON-MOTORIZED table screenshot, found none"
        )
    if len(pmi_chart_pics) < 2:
        logger.warning(
            "Slide 11 PMI expected 2 Excel graph screenshots (Motorized/Stationary), found %s",
            len(pmi_chart_pics),
        )
    if pmi_ole:
        logger.warning("Slide 11 PMI still has %s embedded OLE object(s)", pmi_ole)
    if pmi_charts:
        logger.warning(
            "Slide 11 PMI still has %s native chart(s) — expected Excel graph screenshots",
            pmi_charts,
        )
    pmi_bodies = [
        s
        for s in pmi.shapes
        if getattr(s, "has_text_frame", False)
        and "textbox" in (s.name or "").casefold()
        and int(s.height) > 350_000
        and int(s.top) > 1_000_000
        and int(s.left) >= 8_000_000
    ]
    pmi_leftover = [
        (s.text_frame.text or "").strip()
        for s in pmi_bodies
        if (s.text_frame.text or "").strip()
        and "leading" not in (s.text_frame.text or "").casefold()
        and not (s.text_frame.text or "").casefold().startswith("action plan")
    ]
    if pmi_leftover:
        logger.warning("Slide 11 PMI narrative still has text: %s", pmi_leftover[:2])
    else:
        print("VERIFY slide 11 PMI narrative: Leading Issues / Action Plan boxes are empty")

    # Slide 12 ISR: Regions Rel/Sev table + 2 Excel graph screenshots; narrative empty.
    isr = prs.slides[11]
    isr_charts = sum(1 for s in isr.shapes if getattr(s, "has_chart", False))
    isr_pics = [
        s
        for s in isr.shapes
        if s.shape_type == MSO_SHAPE_TYPE.PICTURE and int(s.top) < 6_200_000
    ]
    isr_table_pics = [s for s in isr_pics if int(s.top) < 4_000_000]
    isr_chart_pics = [s for s in isr_pics if int(s.top) >= 4_000_000]
    print(
        f"VERIFY slide 12 ISR: {len(isr_table_pics)} table pic(s), "
        f"{len(isr_chart_pics)} graph screenshot(s), {isr_charts} native chart(s)"
    )
    if len(isr_table_pics) < 1:
        logger.warning(
            "Slide 12 ISR expected Workings Regions RELIABILITY/SEVERITY table screenshot, found none"
        )
    if len(isr_chart_pics) < 2:
        logger.warning(
            "Slide 12 ISR expected 2 Excel graph screenshots (Reliability/Severity), found %s",
            len(isr_chart_pics),
        )
    if isr_charts:
        logger.warning(
            "Slide 12 ISR still has %s native chart(s) — expected Excel graph screenshots",
            isr_charts,
        )
    isr_bodies = [
        s
        for s in isr.shapes
        if getattr(s, "has_text_frame", False)
        and "textbox" in (s.name or "").casefold()
        and int(s.height) > 350_000
        and int(s.top) > 1_000_000
        and int(s.left) >= 8_000_000
    ]
    isr_leftover = [
        (s.text_frame.text or "").strip()
        for s in isr_bodies
        if (s.text_frame.text or "").strip()
        and "leading" not in (s.text_frame.text or "").casefold()
        and not (s.text_frame.text or "").casefold().startswith("action plan")
    ]
    if isr_leftover:
        logger.warning("Slide 12 ISR narrative still has text: %s", isr_leftover[:2])
    else:
        print("VERIFY slide 12 ISR narrative: Leading Issues / Action Plan boxes are empty")

    # Slide 13 ISR comments: empty narrative boxes present.
    isr_c = prs.slides[12]
    isr_c_bodies = [
        s
        for s in isr_c.shapes
        if getattr(s, "has_text_frame", False)
        and "textbox" in (s.name or "").casefold()
        and int(s.height) > 500_000
        and int(s.top) > 1_000_000
    ]
    isr_c_leftover = [
        (s.text_frame.text or "").strip()
        for s in isr_c_bodies
        if (s.text_frame.text or "").strip()
        and "leading" not in (s.text_frame.text or "").casefold()
        and not (s.text_frame.text or "").casefold().startswith("action plan")
        and "budget" not in (s.text_frame.text or "").casefold()
    ]
    print(f"VERIFY slide 13 ISR comments body text boxes: {len(isr_c_bodies)}")
    if len(isr_c_bodies) < 2:
        logger.warning(
            "Slide 13 ISR comments expected empty editable narrative text boxes, found %s",
            len(isr_c_bodies),
        )
    if isr_c_leftover:
        logger.warning("Slide 13 ISR comments narrative still has text: %s", isr_c_leftover[:2])
    else:
        print("VERIFY slide 13 ISR comments: Leading Issues / Action Plan boxes are empty")


def build_presentation(data: MprData, config: dict, base_dir: Path) -> Path:
    ppt_cfg = config["powerpoint"]
    month_name = date(data.year, data.month, 1).strftime("%B")
    output_name = ppt_cfg.get("output_name", "GSE MPR - {month_name} {year}.pptx").format(
        year=data.year,
        month=data.month,
        month_name=month_name,
        report_title=data.report_output_title(),
    )
    output_dir = base_dir / ppt_cfg["output_dir"]
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / output_name

    prs = _load_template(config, base_dir)

    slide_specs = load_template_map(base_dir)
    if slide_specs:
        spec_by_index = {spec["index"]: spec for spec in slide_specs}
        for idx, slide in enumerate(prs.slides):
            spec = spec_by_index.get(idx)
            if spec:
                _apply_slide_spec(slide, data, config, spec)
                logger.info("Processed slide %s (%s)", idx, spec.get("name", ""))
            elif ppt_cfg.get("clear_unmapped_slides", False):
                _clear_non_title_content(slide)
    else:
        logger.warning("template_map.yaml not found — applying month tokens on slide 0 only.")
        if prs.slides:
            replace_month_tokens_on_slide(prs.slides[0], data)

    _verify_output(prs)
    prs.save(output_path)
    logger.info("Saved report: %s", output_path)
    print(f"\n>>> Open this file to review changes:\n>>> {output_path}\n")
    return output_path
