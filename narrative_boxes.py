"""Clear Leading Issues / Action Plan narrative areas to blank textboxes."""

from __future__ import annotations

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ppt_format import clear_text_frame_content, set_text_frame_preserve

HEADER_LABELS = {
    "leading issues",
    "action plan",
    "action plans",
}


def _iter_all_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_all_shapes(shape.shapes)


def _header_only_text(first_line: str) -> str | None:
    """Return the header label when the first line is a narrative section title."""
    text = first_line.strip()
    lower = text.lower().rstrip(":")
    if lower in HEADER_LABELS:
        return text.rstrip(":").strip()
    if lower.startswith("leading issues"):
        return "Leading Issues"
    if lower.startswith("action plan"):
        return "Action Plans" if lower.startswith("action plans") else "Action Plan"
    return None


def clear_manual_narrative_boxes(slide) -> None:
    """
    Leave Leading Issues / Action Plan areas as blank textboxes.

    Header-only shapes keep their title; mixed header/body shapes keep the
    header line and drop narrative body text.
    """
    for shape in _iter_all_shapes(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        lines = [line.strip() for line in text.splitlines() if line.strip()]
        if not lines:
            continue

        header = _header_only_text(lines[0])
        if header and len(lines) == 1:
            set_text_frame_preserve(shape.text_frame, header)
            continue

        lower = text.lower()
        if header and len(lines) > 1:
            set_text_frame_preserve(shape.text_frame, header)
            continue

        if any(token in lower for token in ("leading issues", "action plan")):
            if header:
                set_text_frame_preserve(shape.text_frame, header)
            else:
                clear_text_frame_content(shape.text_frame)
